#!/usr/bin/env python3
"""
html2pptx.py - Convert Amplifier Stories HTML decks to PowerPoint presentations.

This tool parses HTML presentation decks created with the Amplifier Stories style
and generates equivalent PowerPoint presentations using python-pptx.

Usage:
    uv run --with python-pptx,beautifulsoup4,lxml python tools/html2pptx.py <input.html> [output.pptx]

If output path is not specified, uses the input filename with .pptx extension.

Supports all Amplifier Stories element types including:
- Slide structure (div.slide and section.slide)
- Headlines, subheads, section labels
- Cards in grid layouts (thirds, halves, fourths, grid-2 through grid-5)
- Code blocks with syntax highlighting
- Flow diagrams and architecture diagrams
- Stats grids and stat rows
- Tenet boxes, highlight boxes
- Feature lists, tables, versus comparisons
- Notification stacks
- Comparison tables
- CSS variable extraction for per-deck theming
"""

import argparse
import math
import re
import sys
import warnings
from copy import deepcopy
from pathlib import Path
from typing import Optional

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Slide dimensions ──────────────────────────────────────────────────────────
SLIDE_WIDTH = 10.0  # inches
SLIDE_HEIGHT = 5.625  # inches (16:9)
CONTENT_LEFT = 0.8
CONTENT_WIDTH = 8.4
CONTENT_RIGHT = CONTENT_LEFT + CONTENT_WIDTH

# ── Spacing constants (consistent gap system) ────────────────────────────────
GAP_TIGHT = 0.08  # within a card, between tightly related elements
GAP_NORMAL = 0.10  # standard gap between elements on a slide
GAP_SECTION = 0.20  # gap between major sections (e.g., headline to cards)

# ── Color palette (matching Amplifier Stories dark-mode style) ────────────────
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
MS_CYAN = RGBColor(0x50, 0xE6, 0xFF)
MS_GREEN = RGBColor(0x00, 0xCC, 0x6A)
MS_ORANGE = RGBColor(0xFF, 0x9F, 0x0A)
MS_RED = RGBColor(0xFF, 0x45, 0x3A)
MS_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GRAY_70 = RGBColor(0xB3, 0xB3, 0xB3)
GRAY_50 = RGBColor(0x80, 0x80, 0x80)
DARK_GRAY = RGBColor(0x1A, 0x1A, 0x1A)
BORDER_GRAY = RGBColor(0x33, 0x33, 0x33)
CODE_BG = RGBColor(0x0D, 0x11, 0x17)

# Code syntax colors
CODE_GREEN = RGBColor(0x4A, 0xDE, 0x80)
CODE_BLUE = RGBColor(0x60, 0xA5, 0xFA)
CODE_YELLOW = RGBColor(0xFB, 0xBF, 0x24)
CODE_GRAY = RGBColor(0x6B, 0x73, 0x80)
CODE_PURPLE = RGBColor(0xC0, 0x84, 0xFC)
CODE_STRING = RGBColor(0xFB, 0xBF, 0x24)
CODE_DEFAULT = RGBColor(0xE6, 0xE6, 0xE6)

# ── Default font ──────────────────────────────────────────────────────────────
DEFAULT_FONT = "Arial"
CODE_FONT = "Consolas"


# ── Helper: hex string → RGBColor ────────────────────────────────────────────
def hex_to_rgb(hex_str: str) -> Optional[RGBColor]:
    """Convert #RGB or #RRGGBB to RGBColor."""
    hex_str = hex_str.strip().lstrip("#")
    if len(hex_str) == 3:
        hex_str = "".join(c * 2 for c in hex_str)
    if len(hex_str) != 6:
        return None
    try:
        r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        return None


def parse_color_from_class(classes: list[str]) -> Optional[RGBColor]:
    """Extract accent color from CSS classes."""
    color_map = {
        "green": MS_GREEN,
        "orange": MS_ORANGE,
        "red": MS_RED,
        "ms-green": MS_GREEN,
        "ms-orange": MS_ORANGE,
        "ms-red": MS_RED,
        "ms-blue": MS_BLUE,
        "ms-cyan": MS_CYAN,
        "ms-purple": MS_PURPLE,
        "warning": MS_ORANGE,
    }
    for cls in classes:
        if cls in color_map:
            return color_map[cls]
    return None


def _replace_br_tags(element: Tag):
    """Replace all <br> tags with newline characters BEFORE text extraction."""
    for br in element.find_all("br"):
        br.replace_with("\n")


def get_text(element: Optional[Tag]) -> str:
    """Extract text content from an element, handling None and <br> tags.

    Uses a space separator so inline elements (<code>, <span>, etc.)
    don't get their text fused with surrounding words.
    """
    if element is None:
        return ""
    # Work on a copy so we don't mutate the original soup
    el_copy = deepcopy(element)
    _replace_br_tags(el_copy)
    # Use " " separator to prevent word-joining across inline tags
    text = el_copy.get_text(" ")
    # Collapse runs of spaces (but keep newlines)
    lines = text.split("\n")
    lines = [" ".join(line.split()) for line in lines]
    return "\n".join(lines).strip()


def get_rich_text(element: Optional[Tag]) -> list[dict]:
    """Extract rich text runs preserving bold/italic/highlight spans.

    Returns a list of dicts: [{"text": str, "bold": bool, "italic": bool, "color": RGBColor|None}]
    """
    if element is None:
        return []

    el_copy = deepcopy(element)
    _replace_br_tags(el_copy)
    runs: list[dict] = []

    for child in el_copy.descendants:
        if isinstance(child, NavigableString):
            text = str(child)
            if not text:
                continue
            # Collapse HTML indentation whitespace (keep \n from <br> tags)
            lines = text.split("\n")
            lines = [" ".join(line.split()) for line in lines]
            text = "\n".join(lines)
            if not text or text.isspace():
                continue
            # Determine formatting from parent chain
            bold = False
            italic = False
            color = None
            parent = child.parent
            while parent and parent.name:
                if parent.name in ("strong", "b"):
                    bold = True
                if parent.name in ("em", "i"):
                    italic = True
                if parent.name == "span":
                    cls = parent.get("class", [])
                    if "highlight" in cls:
                        color = MS_CYAN
                    elif "check" in cls:
                        color = MS_GREEN
                parent = parent.parent
            runs.append({"text": text, "bold": bold, "italic": italic, "color": color})

    # Merge adjacent runs with identical formatting
    merged: list[dict] = []
    for run in runs:
        if (
            merged
            and merged[-1]["bold"] == run["bold"]
            and merged[-1]["italic"] == run["italic"]
            and merged[-1]["color"] == run["color"]
        ):
            merged[-1]["text"] += run["text"]
        else:
            merged.append(run)

    # Strip leading whitespace from first run, trailing from last
    if merged:
        merged[0]["text"] = merged[0]["text"].lstrip()
        merged[-1]["text"] = merged[-1]["text"].rstrip()
        # Remove empty runs
        merged = [r for r in merged if r["text"]]

    return merged


# ── Low-level shape helpers ───────────────────────────────────────────────────


def set_slide_background(slide, color=BLACK):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_font(
    run,
    name: str = DEFAULT_FONT,
    size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = WHITE,
):
    """Apply font properties to a text run."""
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    font_name: str = DEFAULT_FONT,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    """Add a text box with specified styling. Every run gets explicit font name."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    tf.vertical_anchor = vertical_anchor

    # Handle multi-line text
    lines = text.split("\n")
    for line_idx, line in enumerate(lines):
        if line_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        _set_font(
            run, name=font_name, size=font_size, bold=bold, italic=italic, color=color
        )

    return box


def add_rich_text_box(
    slide,
    runs: list[dict],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    font_name: str = DEFAULT_FONT,
    default_color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    """Add a text box with multiple formatting runs."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True

    # Split runs by newlines into paragraphs
    paragraphs: list[list[dict]] = [[]]
    for r in runs:
        parts = r["text"].split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                paragraphs.append([])
            if part:
                paragraphs[-1].append({**r, "text": part})

    for p_idx, p_runs in enumerate(paragraphs):
        if p_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if not p_runs:
            # Empty paragraph (blank line)
            run = p.add_run()
            run.text = ""
            _set_font(run, name=font_name, size=font_size, color=default_color)
        else:
            for r in p_runs:
                run = p.add_run()
                run.text = r["text"]
                _set_font(
                    run,
                    name=font_name,
                    size=font_size,
                    bold=r.get("bold", False),
                    italic=r.get("italic", False),
                    color=r.get("color") or default_color,
                )

    return box


def add_filled_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: RGBColor = DARK_GRAY,
    border_color: Optional[RGBColor] = BORDER_GRAY,
    border_width: float = 1.0,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    """Add a filled shape (card background, code block bg, etc.)."""
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


# ── Element-level helpers ─────────────────────────────────────────────────────


# ---------------------------------------------------------------------------
# Arial character-width tables (fraction of em-size)
# Measured from Arial TrueType metrics.  Missing chars use the default.
# ---------------------------------------------------------------------------
_ARIAL_REGULAR: dict[str, float] = {
    " ": 0.28,
    "!": 0.28,
    '"': 0.35,
    "#": 0.56,
    "$": 0.56,
    "%": 0.89,
    "&": 0.67,
    "'": 0.19,
    "(": 0.33,
    ")": 0.33,
    "*": 0.39,
    "+": 0.58,
    ",": 0.28,
    "-": 0.33,
    ".": 0.28,
    "/": 0.28,
    "0": 0.56,
    "1": 0.56,
    "2": 0.56,
    "3": 0.56,
    "4": 0.56,
    "5": 0.56,
    "6": 0.56,
    "7": 0.56,
    "8": 0.56,
    "9": 0.56,
    ":": 0.28,
    ";": 0.28,
    "<": 0.58,
    "=": 0.58,
    ">": 0.58,
    "?": 0.56,
    "@": 1.02,
    "A": 0.67,
    "B": 0.67,
    "C": 0.72,
    "D": 0.72,
    "E": 0.67,
    "F": 0.61,
    "G": 0.78,
    "H": 0.72,
    "I": 0.28,
    "J": 0.50,
    "K": 0.67,
    "L": 0.56,
    "M": 0.83,
    "N": 0.72,
    "O": 0.78,
    "P": 0.67,
    "Q": 0.78,
    "R": 0.72,
    "S": 0.67,
    "T": 0.61,
    "U": 0.72,
    "V": 0.67,
    "W": 0.94,
    "X": 0.67,
    "Y": 0.67,
    "Z": 0.61,
    "[": 0.28,
    "\\": 0.28,
    "]": 0.28,
    "^": 0.47,
    "_": 0.56,
    "`": 0.33,
    "a": 0.56,
    "b": 0.56,
    "c": 0.50,
    "d": 0.56,
    "e": 0.56,
    "f": 0.28,
    "g": 0.56,
    "h": 0.56,
    "i": 0.22,
    "j": 0.22,
    "k": 0.50,
    "l": 0.22,
    "m": 0.83,
    "n": 0.56,
    "o": 0.56,
    "p": 0.56,
    "q": 0.56,
    "r": 0.33,
    "s": 0.50,
    "t": 0.28,
    "u": 0.56,
    "v": 0.50,
    "w": 0.72,
    "x": 0.50,
    "y": 0.50,
    "z": 0.50,
    "{": 0.33,
    "|": 0.26,
    "}": 0.33,
    "~": 0.58,
}
_ARIAL_REGULAR_DEFAULT = 0.56

# Bold glyphs are ~8% wider on average
_ARIAL_BOLD_SCALE = 1.08


def _estimate_text_width_pt(text: str, font_size_pt: int, bold: bool = False) -> float:
    """Return estimated rendered width of *text* in points for Arial."""
    table = _ARIAL_REGULAR
    default = _ARIAL_REGULAR_DEFAULT
    total = 0.0
    for ch in text:
        total += table.get(ch, default)
    width_pt = total * font_size_pt
    if bold:
        width_pt *= _ARIAL_BOLD_SCALE
    return width_pt


def _estimate_text_height(
    text: str,
    font_size_pt: int,
    box_width_inches: float,
    line_spacing: float = 1.2,
    bold: bool = False,
) -> float:
    """Estimate rendered text height using per-character width tables.

    Uses Arial TrueType character-width metrics to compute the actual
    rendered width of each line, then determines wrapping based on the
    available box width.  This is far more accurate than average-factor
    approaches because character widths in proportional fonts vary by 4x
    (e.g. 'i' = 0.22em vs 'W' = 0.94em).

    PowerPoint text frames have ~0.1" internal margin on each side, so the
    usable text width is reduced by 0.2".  The 0.10" vertical padding
    accounts for top+bottom text-frame margins.
    """
    # Usable text width inside the text frame (0.1" margin each side)
    usable_width_pt = (box_width_inches - 0.20) * 72.0
    if usable_width_pt < 36:  # absolute minimum
        usable_width_pt = 36.0

    paragraphs = text.split("\n")
    num_lines = 0.0
    for para in paragraphs:
        stripped = para.strip()
        if not stripped:
            num_lines += 0.4  # empty paragraph gets ~40% of a line
            continue
        # Compute rendered width and determine how many lines it wraps to
        rendered_pt = _estimate_text_width_pt(stripped, font_size_pt, bold)
        if rendered_pt <= usable_width_pt:
            num_lines += 1
        else:
            # Word-wrap estimation: divide rendered width by usable width.
            # Add 5% for word-boundary inefficiency (lines can't break mid-word
            # optimally, so there's leftover space at the end of each line).
            wrap_lines = rendered_pt / usable_width_pt * 1.05
            num_lines += max(2, math.ceil(wrap_lines))

    line_height = font_size_pt / 72.0 * line_spacing
    return num_lines * line_height + 0.10  # 0.10" for text-frame margins


def _truncate_to_fit(
    text: str,
    font_size_pt: int,
    box_width_inches: float,
    max_height_inches: float,
    bold: bool = False,
) -> str:
    """Truncate text so it fits within max_height_inches at the given font size.

    Uses binary search on paragraph boundaries to find the longest prefix
    that fits, appending '...' when truncation is needed.
    """
    if (
        _estimate_text_height(text, font_size_pt, box_width_inches, bold=bold)
        <= max_height_inches
    ):
        return text  # already fits

    paragraphs = text.split("\n")
    # Binary search for how many paragraphs fit
    lo, hi = 0, len(paragraphs)
    while lo < hi:
        mid = (lo + hi + 1) // 2
        candidate = "\n".join(paragraphs[:mid]) + "\n..."
        if (
            _estimate_text_height(candidate, font_size_pt, box_width_inches, bold=bold)
            <= max_height_inches
        ):
            lo = mid
        else:
            hi = mid - 1
    if lo > 0:
        return "\n".join(paragraphs[:lo]) + "\n..."

    # Single long paragraph — truncate by characters
    words = text.split()
    lo, hi = 0, len(words)
    while lo < hi:
        mid = (lo + hi + 1) // 2
        candidate = " ".join(words[:mid]) + "..."
        if (
            _estimate_text_height(candidate, font_size_pt, box_width_inches, bold=bold)
            <= max_height_inches
        ):
            lo = mid
        else:
            hi = mid - 1
    return " ".join(words[: max(lo, 1)]) + "..."


def add_section_label(slide, text: str, top: float = 0.6, color: RGBColor = MS_BLUE):
    """Add a colored uppercase section label."""
    return add_text_box(
        slide,
        text.upper(),
        left=CONTENT_LEFT,
        top=top,
        width=CONTENT_WIDTH,
        height=0.4,
        font_size=14,
        bold=True,
        color=color,
        wrap=False,
    )


def add_headline(
    slide,
    text: str,
    top: float = 1.1,
    size: int = 48,
    center: bool = False,
    color: RGBColor = WHITE,
):
    """Add a large headline with dynamically estimated box height."""
    # Estimate actual height needed — headlines are always bold
    box_height = _estimate_text_height(text, size, CONTENT_WIDTH, bold=True)
    # At least one line; cap at 3.5" (overflow compression handles the rest)
    min_height = size / 72 * 1.2 + 0.1  # single line minimum
    box_height = max(min_height, min(box_height, 3.5))
    return add_text_box(
        slide,
        text,
        left=CONTENT_LEFT,
        top=top,
        width=CONTENT_WIDTH,
        height=box_height,
        font_size=size,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT,
    )


def add_subhead(
    slide, text: str, top: float = 2.8, color: RGBColor = GRAY_70, center: bool = False
):
    """Add a subtitle/subheading with dynamically estimated box height."""
    # Estimate height needed — subheads are not bold
    box_height = _estimate_text_height(text, 24, CONTENT_WIDTH, bold=False)
    min_height = 24 / 72 * 1.2 + 0.1  # single line minimum
    box_height = max(min_height, min(box_height, 2.5))
    return add_text_box(
        slide,
        text,
        left=CONTENT_LEFT,
        top=top,
        width=CONTENT_WIDTH,
        height=box_height,
        font_size=24,
        color=color,
        align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT,
    )


def add_card(
    slide,
    title: str,
    text: str,
    left: float,
    top: float,
    width: float = 2.6,
    height: float = 1.8,
    title_color: RGBColor = MS_BLUE,
    rich_runs: Optional[list[dict]] = None,
):
    """Add a card with title and description."""
    # Card background
    add_filled_box(slide, left, top, width, height)

    # Adaptive title height: estimate from text, allow up to 50% of card.
    # The +0.05" buffer prevents tight-fit overflows where text-frame internal
    # margins cause the estimate to be fractionally too small.
    title_est = _estimate_text_height(title, 16, width - 0.3, bold=True) + 0.05
    title_height = max(0.35, min(title_est, height * 0.50))

    # Card title
    add_text_box(
        slide,
        title,
        left=left + 0.15,
        top=top + 0.15,
        width=width - 0.3,
        height=title_height,
        font_size=16,
        bold=True,
        color=title_color,
    )

    # Card text (rich or plain)
    text_top = top + 0.15 + title_height + GAP_TIGHT  # title_top + title_height + gap
    text_height = height - (
        0.15 + title_height + GAP_TIGHT + 0.05
    )  # pad + title + gap + bottom_pad

    # Auto-reduce font when description overflows available space
    inner_w = width - 0.3
    desc_content = text if text else ""
    if rich_runs:
        desc_content = " ".join(r.get("text", "") for r in rich_runs)
    desc_font = 12
    for try_size in (12, 11, 10, 9, 8):
        if _estimate_text_height(desc_content, try_size, inner_w) <= max(
            text_height, 0.2
        ):
            desc_font = try_size
            break
        desc_font = try_size  # use smallest if none fit

    # Ensure text box is tall enough for the text at the chosen font size
    desc_est = _estimate_text_height(desc_content, desc_font, inner_w)
    text_height = max(text_height, desc_est)

    if rich_runs:
        add_rich_text_box(
            slide,
            rich_runs,
            left=left + 0.15,
            top=text_top,
            width=inner_w,
            height=max(text_height, 0.2),
            font_size=desc_font,
            default_color=GRAY_70,
        )
    else:
        add_text_box(
            slide,
            text,
            left=left + 0.15,
            top=text_top,
            width=inner_w,
            height=max(text_height, 0.2),
            font_size=desc_font,
            color=GRAY_70,
        )


def add_tenet(
    slide,
    title: str,
    text: str,
    left: float,
    top: float,
    width: float = 4.0,
    height: float = 0.9,
    accent_color: RGBColor = MS_GREEN,
):
    """Add a tenet box with left border accent."""
    bg_map = {
        MS_GREEN: RGBColor(0x0D, 0x1A, 0x0D),
        MS_ORANGE: RGBColor(0x1A, 0x15, 0x0D),
        MS_RED: RGBColor(0x1A, 0x0D, 0x0D),
    }
    bg_color = bg_map.get(accent_color, RGBColor(0x0D, 0x15, 0x1A))

    # Background
    add_filled_box(
        slide, left, top, width, height, fill_color=bg_color, border_color=None
    )

    # Left accent bar (0.15" wide for visibility in PowerPoint)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.15), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Remove text frame to keep it a pure colored rectangle
    bar.text_frame.word_wrap = False

    # Title (offset right to clear the 0.15" accent bar)
    text_left = left + 0.25
    text_width = width - 0.35
    title_h = _estimate_text_height(title, 14, text_width, bold=True)
    title_h = max(0.25, min(title_h, height * 0.45))
    add_text_box(
        slide,
        title,
        left=text_left,
        top=top + 0.08,
        width=text_width,
        height=title_h,
        font_size=14,
        bold=True,
        color=WHITE,
    )

    # Text
    text_top = top + 0.08 + title_h + GAP_TIGHT
    text_h = height - (0.08 + title_h + GAP_TIGHT + 0.05)
    if text and text_h > 0.1:
        add_text_box(
            slide,
            text,
            left=text_left,
            top=text_top,
            width=text_width,
            height=max(text_h, 0.15),
            font_size=11,
            color=GRAY_70,
        )


def add_highlight_box(
    slide,
    text: str,
    top: float = 4.2,
    color: RGBColor = MS_BLUE,
    rich_runs: Optional[list[dict]] = None,
):
    """Add a highlight/callout box."""
    bg_map = {
        MS_GREEN: RGBColor(0x00, 0x1A, 0x0D),
        MS_ORANGE: RGBColor(0x33, 0x1A, 0x00),
        MS_RED: RGBColor(0x1A, 0x0D, 0x0D),
    }
    bg_color = bg_map.get(color, RGBColor(0x00, 0x1A, 0x33))

    # Estimate height from content (use bold=True since highlight boxes
    # commonly contain <strong> emphasis which widens character metrics)
    inner_w = CONTENT_WIDTH - 0.4
    has_bold = rich_runs and any(r.get("bold") for r in rich_runs)
    text_h = _estimate_text_height(text, 14, inner_w, bold=has_bold or False)
    box_h = max(0.5, text_h + 0.24)  # 0.12 pad top + bottom

    add_filled_box(
        slide,
        CONTENT_LEFT,
        top,
        CONTENT_WIDTH,
        box_h,
        fill_color=bg_color,
        border_color=color,
        border_width=1,
    )

    if rich_runs:
        add_rich_text_box(
            slide,
            rich_runs,
            left=CONTENT_LEFT + 0.2,
            top=top + 0.12,
            width=inner_w,
            height=box_h - 0.24,
            font_size=14,
            default_color=WHITE,
        )
    else:
        add_text_box(
            slide,
            text,
            left=CONTENT_LEFT + 0.2,
            top=top + 0.12,
            width=inner_w,
            height=box_h - 0.24,
            font_size=14,
            color=WHITE,
        )


# ── CSS variable extraction ──────────────────────────────────────────────────


def extract_css_vars(soup: BeautifulSoup) -> dict[str, str]:
    """Extract CSS custom properties from <style> blocks."""
    result: dict[str, str] = {}
    for style_tag in soup.find_all("style"):
        css_text = style_tag.string or ""
        # Match --variable-name: value patterns
        for m in re.finditer(r"--([\w-]+)\s*:\s*([^;]+)", css_text):
            name = m.group(1).strip()
            value = m.group(2).strip()
            result[name] = value
    return result


def resolve_accent_color(css_vars: dict[str, str]) -> RGBColor:
    """Determine the accent color for the deck from CSS variables."""
    # Try common variable names
    for var_name in ("color-accent", "accent"):
        val = css_vars.get(var_name, "")
        if val.startswith("#"):
            rgb = hex_to_rgb(val)
            if rgb:
                return rgb
    return MS_BLUE


# ── Main converter class ─────────────────────────────────────────────────────


class HTMLToPPTXConverter:
    """Converts Amplifier Stories HTML decks to PowerPoint."""

    def __init__(self, html_content: str):
        self.soup = BeautifulSoup(html_content, "lxml")
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_WIDTH)
        self.prs.slide_height = Inches(SLIDE_HEIGHT)
        self.blank_layout = self.prs.slide_layouts[6]

        # Extract per-deck theming from CSS variables
        self.css_vars = extract_css_vars(self.soup)
        self.accent_color = resolve_accent_color(self.css_vars)

        # Track warnings
        self.warnings: list[str] = []

    # ── Slide extraction ─────────────────────────────────────────────────────

    def extract_slides(self) -> list[Tag]:
        """Extract all slide elements (div.slide AND section.slide)."""
        return self.soup.find_all(["div", "section"], class_="slide")

    def is_centered(self, slide_div: Tag) -> bool:
        """Check if slide has center or title-slide class."""
        classes = slide_div.get("class", [])
        return "center" in classes or "title-slide" in classes

    # ── Main slide processor ─────────────────────────────────────────────────

    def process_slide(self, slide_div: Tag, slide_num: int):
        """Process a single slide and add to presentation."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        set_slide_background(slide)

        is_centered = self.is_centered(slide_div)
        current_top = 0.6

        # Track which child elements we've handled (to avoid double-processing)
        handled_elements: set[int] = set()

        # ── Section number + section title (alternative to section-label) ──
        section_number = slide_div.find(class_="section-number")
        if section_number:
            add_section_label(
                slide,
                get_text(section_number),
                top=current_top,
                color=self.accent_color,
            )
            current_top += 0.4 + GAP_TIGHT  # section-label height + tight gap
            handled_elements.add(id(section_number))

        section_title = slide_div.find(class_="section-title")
        if section_title:
            st_text = get_text(section_title)
            st_height = _estimate_text_height(st_text, 20, CONTENT_WIDTH, bold=True)
            st_height = max(0.40, min(st_height, 1.0))
            add_text_box(
                slide,
                st_text,
                left=CONTENT_LEFT,
                top=current_top,
                width=CONTENT_WIDTH,
                height=st_height,
                font_size=20,
                bold=True,
                color=WHITE,
            )
            current_top += st_height + GAP_NORMAL
            handled_elements.add(id(section_title))

        # ── Section label ────────────────────────────────────────────────────
        section_label = slide_div.find(class_="section-label")
        if section_label:
            if is_centered:
                current_top = 1.0
            add_section_label(
                slide, get_text(section_label), top=current_top, color=self.accent_color
            )
            current_top += 0.4 + GAP_NORMAL  # section-label height + normal gap
            handled_elements.add(id(section_label))

        # ── Headline (h1 or h2.headline) ─────────────────────────────────────
        headline = slide_div.find(["h1", "h2"], class_="headline") or slide_div.find(
            "h1"
        )
        if headline:
            text = get_text(headline)
            has_gradient = "big-text" in headline.get("class", [])
            color = MS_CYAN if has_gradient else WHITE
            size = 56 if headline.name == "h1" or has_gradient else 40

            if is_centered:
                current_top = max(current_top, 1.5)

            add_headline(
                slide, text, top=current_top, size=size, center=is_centered, color=color
            )
            # Advance by dynamically estimated headline height + section gap
            headline_height = _estimate_text_height(
                text, size, CONTENT_WIDTH, bold=True
            )
            min_h = size / 72 * 1.2 + 0.1
            headline_height = max(min_h, min(headline_height, 3.5))
            # Fixed 0.50" gap on centered/title slides for clean title-to-subtitle
            # spacing; regular slides use standard section gap.
            gap_after = 0.50 if is_centered else GAP_SECTION
            current_top += headline_height + gap_after
            handled_elements.add(id(headline))

        # ── Medium headline ──────────────────────────────────────────────────
        medium_headline = slide_div.find(class_="medium-headline")
        if medium_headline and id(medium_headline) not in handled_elements:
            mh_text = get_text(medium_headline)
            add_headline(
                slide,
                mh_text,
                top=current_top,
                size=36,
                center=is_centered,
            )
            mh_height = _estimate_text_height(mh_text, 36, CONTENT_WIDTH, bold=True)
            mh_height = max(36 / 72 * 1.2 + 0.1, min(mh_height, 3.0))
            current_top += mh_height + GAP_SECTION
            handled_elements.add(id(medium_headline))

        # ── Subhead ──────────────────────────────────────────────────────────
        subhead = slide_div.find(class_="subhead")
        if subhead:
            text = get_text(subhead)
            add_subhead(slide, text, top=current_top, center=is_centered)
            # Advance by the same box height that add_subhead() computed + gap
            sub_height = _estimate_text_height(text, 24, CONTENT_WIDTH, bold=False)
            min_sub = 24 / 72 * 1.2 + 0.1
            sub_height = max(
                min_sub, min(sub_height, 2.5)
            )  # must match add_subhead caps
            current_top += sub_height + GAP_NORMAL
            handled_elements.add(id(subhead))

        # ── Architecture diagram (preformatted text art) ─────────────────────
        arch_diagram = slide_div.find(class_="architecture-diagram")
        if arch_diagram:
            self._add_architecture_diagram(slide, arch_diagram, current_top)
            current_top += 3.0
            handled_elements.add(id(arch_diagram))

        # ── Comparison table (grid layout or HTML table) ──────────────────────
        comp_table = slide_div.find(class_="comparison-table")
        if comp_table:
            if comp_table.name == "table":
                # It's a real HTML <table> with class comparison-table —
                # route to the table handler instead of the CSS-grid handler.
                self._add_table(slide, comp_table, current_top)
                rows = comp_table.find_all("tr")
                current_top += len(rows) * 0.32 + 0.3
            else:
                self._add_comparison_table(slide, comp_table, current_top)
                current_top += 2.2
            handled_elements.add(id(comp_table))

        # ── Cards in grid containers ─────────────────────────────────────────
        grid_classes = [
            "thirds",
            "halves",
            "fourths",
            "grid",
            "grid-2",
            "grid-3",
            "grid-4",
            "grid-5",
            "tools-grid",
        ]
        # Exclude stat/velocity grids (handled by _add_stats) and principles
        # NOTE: BS4 class_=lambda receives individual class strings, not lists
        _grid_exclude = {"principles-grid", "stat-grid", "velocity-grid"}
        card_containers = slide_div.find_all(
            class_=lambda c: (
                c and any(gc in c for gc in grid_classes) and c not in _grid_exclude
            )
        )
        for container in card_containers:
            if id(container) in handled_elements:
                continue

            # Find cards or module-cards inside
            cards = container.find_all(
                class_=["card", "module-card", "tool-card"], recursive=False
            )
            if not cards:
                # Try deeper (some decks nest differently)
                cards = container.find_all(class_=["card", "module-card", "tool-card"])

            # Check for code-blocks inside grid (e.g. protocol definitions)
            inner_code_blocks = container.find_all(class_="code-block")

            if cards:
                cards_height = self._add_cards(slide, cards, current_top, container)
                current_top += cards_height + GAP_SECTION
                for c in cards:
                    handled_elements.add(id(c))

                # Also process non-card siblings in the grid that have content
                direct_children = [
                    ch
                    for ch in container.children
                    if isinstance(ch, Tag) and id(ch) not in handled_elements
                ]
                for child in direct_children:
                    child_classes = child.get("class", [])
                    # Skip if it's a card (already handled)
                    if any(
                        cc in child_classes
                        for cc in ["card", "module-card", "tool-card"]
                    ):
                        continue
                    # Render non-card grid children as generic content
                    child_text = get_text(child)
                    if child_text and len(child_text) > 5:
                        child_text = _truncate_to_fit(
                            child_text, 12, CONTENT_WIDTH, 5.5
                        )
                        child_h = _estimate_text_height(child_text, 12, CONTENT_WIDTH)
                        child_h = max(0.3, child_h)
                        add_text_box(
                            slide,
                            child_text,
                            left=CONTENT_LEFT,
                            top=current_top,
                            width=CONTENT_WIDTH,
                            height=child_h,
                            font_size=12,
                            color=GRAY_70,
                        )
                        current_top += child_h + GAP_NORMAL
                        handled_elements.add(id(child))

                handled_elements.add(id(container))

            elif inner_code_blocks:
                # Grid contains code blocks — render side-by-side if grid-2
                container_cls = " ".join(container.get("class", []))
                num_cb = len(inner_code_blocks)
                if "grid-2" in container_cls and num_cb == 2:
                    # Side-by-side code blocks
                    cb_gap = 0.15
                    cb_width = (CONTENT_WIDTH - cb_gap) / 2
                    max_lines = 0
                    for col_idx, cb in enumerate(inner_code_blocks):
                        cb_left = CONTENT_LEFT + col_idx * (cb_width + cb_gap)
                        self._add_code_block_sized(
                            slide, cb, current_top, cb_left, cb_width
                        )
                        lines = get_text(cb).count("\n") + 1
                        max_lines = max(max_lines, lines)
                        handled_elements.add(id(cb))
                    current_top += max(1.2, min(max_lines * 0.18 + 0.4, 3.5))
                else:
                    # Stacked code blocks (default)
                    for cb in inner_code_blocks:
                        self._add_code_block(slide, cb, current_top)
                        lines = get_text(cb).count("\n") + 1
                        current_top += max(1.2, min(lines * 0.18 + 0.4, 3.5))
                        handled_elements.add(id(cb))
                handled_elements.add(id(container))

            elif container.find(class_="principle"):
                # Grid contains principles — let the principle handler below
                # pick them up; do NOT mark container as handled.
                pass

            else:
                # Grid with unrecognized content — render children as text
                direct_children = [
                    ch
                    for ch in container.children
                    if isinstance(ch, Tag)
                    and ch.get_text(strip=True)
                    and id(ch) not in handled_elements
                ]
                if direct_children:
                    for child in direct_children:
                        child_text = get_text(child)
                        if child_text and len(child_text) > 5:
                            child_text = _truncate_to_fit(
                                child_text, 12, CONTENT_WIDTH, 5.5
                            )
                            child_h = _estimate_text_height(
                                child_text, 12, CONTENT_WIDTH
                            )
                            child_h = max(0.3, child_h)
                            add_text_box(
                                slide,
                                child_text,
                                left=CONTENT_LEFT,
                                top=current_top,
                                width=CONTENT_WIDTH,
                                height=child_h,
                                font_size=12,
                                color=GRAY_70,
                            )
                            current_top += child_h + GAP_NORMAL
                            handled_elements.add(id(child))
                    handled_elements.add(id(container))

        # ── Standalone cards not in containers ───────────────────────────────
        standalone_cards = [
            c
            for c in slide_div.find_all(class_=["card", "module-card", "tool-card"])
            if id(c) not in handled_elements
            and not c.find_parent(
                class_=lambda x: x and any(gc in x for gc in grid_classes)
            )
        ]
        if standalone_cards:
            cards_height = self._add_cards(slide, standalone_cards, current_top)
            current_top += cards_height + GAP_SECTION
            for c in standalone_cards:
                handled_elements.add(id(c))

        # ── Principle boxes (numbered principles in grid layouts) ─────────
        principles = [
            p
            for p in slide_div.find_all(class_="principle")
            if id(p) not in handled_elements
        ]
        if principles:
            p_row_height = self._add_principles(slide, principles, current_top)
            num_rows = -(-len(principles) // 2)  # ceil division
            row_gap = 0.1
            current_top += num_rows * (p_row_height + row_gap)
            for p in principles:
                handled_elements.add(id(p))
            # Also mark the principles-grid container as handled
            pg = slide_div.find(class_="principles-grid")
            if pg:
                handled_elements.add(id(pg))

        # ── Code blocks ──────────────────────────────────────────────────────
        code_blocks = slide_div.find_all(class_="code-block")
        for cb in code_blocks:
            if id(cb) in handled_elements:
                continue
            self._add_code_block(slide, cb, current_top)
            # Estimate height from line count
            lines = get_text(cb).count("\n") + 1
            current_top += max(1.2, min(lines * 0.18 + 0.4, 3.5))
            handled_elements.add(id(cb))

        # ── Flow diagrams ────────────────────────────────────────────────────
        flow_diagrams = slide_div.find_all(class_=["flow-diagram", "workflow", "flow"])
        for fd in flow_diagrams:
            if id(fd) in handled_elements:
                continue
            fd_height = self._add_flow_diagram(slide, fd, current_top)
            current_top += (fd_height or 1.2) + GAP_SECTION
            handled_elements.add(id(fd))

        # ── Tenet boxes ──────────────────────────────────────────────────────
        tenets = slide_div.find_all(class_="tenet")
        if tenets:
            self._add_tenets(slide, tenets, current_top)
            current_top += len(tenets) * 0.5 + 0.5

        # ── Versus comparison ────────────────────────────────────────────────
        versus = slide_div.find(class_="versus")
        if versus:
            self._add_versus(slide, versus, current_top)
            current_top += 2.5

        # ── Tables (data-table or plain table) ───────────────────────────────
        tables = slide_div.find_all("table")
        for table in tables:
            if id(table) in handled_elements:
                continue
            self._add_table(slide, table, current_top)
            rows = table.find_all("tr")
            current_top += len(rows) * 0.32 + 0.3
            handled_elements.add(id(table))

        # ── Feature lists ────────────────────────────────────────────────────
        feature_lists = slide_div.find_all(class_="feature-list")
        for fl in feature_lists:
            if fl.find_parent(class_="versus"):
                continue
            self._add_feature_list(slide, fl, current_top)
            items = fl.find_all("li")
            current_top += len(items) * 0.35 + GAP_SECTION
            handled_elements.add(id(fl))
            for item in items:
                handled_elements.add(id(item))

        # ── Notification stacks ──────────────────────────────────────────────
        notif_stack = slide_div.find(class_="notification-stack")
        if notif_stack and id(notif_stack) not in handled_elements:
            self._add_notification_stack(slide, notif_stack, current_top)
            current_top += 2.5
            handled_elements.add(id(notif_stack))

        # ── Stats grid / stat row / velocity-grid ─────────────────────────────
        stat_container = (
            slide_div.find(class_="stat-grid")
            or slide_div.find(class_="stat-row")
            or slide_div.find(class_="velocity-grid")
        )
        if stat_container and id(stat_container) not in handled_elements:
            # velocity-grid uses velocity-stat/velocity-number/velocity-label
            # which map to stat/stat-number/stat-label
            self._add_stats(slide, stat_container, current_top)
            current_top += 1.2
            handled_elements.add(id(stat_container))

        # ── Big stats (large number + unit, e.g. "90% savings") ──────────────
        big_stats = [
            bs
            for bs in slide_div.find_all(class_="big-stat")
            if id(bs) not in handled_elements
        ]
        if big_stats:
            width_per = CONTENT_WIDTH / len(big_stats)
            for i, bs in enumerate(big_stats):
                num_el = bs.find(class_="big-stat-number")
                unit_el = bs.find(class_="big-stat-unit")
                number = get_text(num_el) if num_el else ""
                unit = get_text(unit_el) if unit_el else ""
                left = CONTENT_LEFT + i * width_per
                num_h = _estimate_text_height(number, 56, width_per, bold=True)
                num_h = max(1.05, num_h)  # 56pt needs ~1.03" minimum
                add_text_box(
                    slide,
                    number,
                    left=left,
                    top=current_top,
                    width=width_per,
                    height=num_h,
                    font_size=56,
                    bold=True,
                    color=MS_CYAN,
                    align=PP_ALIGN.CENTER,
                )
                unit_h = _estimate_text_height(unit, 18, width_per)
                unit_h = max(0.4, unit_h)
                add_text_box(
                    slide,
                    unit,
                    left=left,
                    top=current_top + num_h,
                    width=width_per,
                    height=unit_h,
                    font_size=18,
                    color=GRAY_70,
                    align=PP_ALIGN.CENTER,
                )
                handled_elements.add(id(bs))
            current_top += 1.3 + GAP_SECTION

        # ── Tier stack (tiered info layout) ──────────────────────────────────
        tier_stack = slide_div.find(class_="tier-stack")
        if tier_stack and id(tier_stack) not in handled_elements:
            tiers = tier_stack.find_all(class_="tier")
            num_tiers = len(tiers)
            available = SLIDE_HEIGHT - current_top - 0.4  # bottom margin
            tier_gap = 0.08
            inner_w = CONTENT_WIDTH - 0.3

            # First pass: extract content and compute needed tier height
            tier_items: list[tuple[str, str, str]] = []
            max_tier_need = 0.55  # minimum
            for tier in tiers:
                label_el = tier.find(class_="tier-label")
                title_el = tier.find(class_="tier-title")
                desc_el = tier.find(class_="tier-desc")
                tokens_el = tier.find(class_="tier-tokens")
                label = get_text(label_el) if label_el else ""
                title = get_text(title_el) if title_el else ""
                desc = get_text(desc_el) if desc_el else ""
                tokens = get_text(tokens_el) if tokens_el else ""
                combined = desc
                if tokens:
                    combined = f"{desc}  |  {tokens}" if desc else tokens
                tier_items.append((label, title, combined))
                lh = _estimate_text_height(label.upper(), 10, inner_w, bold=True)
                th = _estimate_text_height(title, 12, inner_w, bold=True)
                dh = _estimate_text_height(combined, 9, inner_w) if combined else 0.0
                need = 0.05 + lh + th + 0.02 + dh + 0.05
                max_tier_need = max(max_tier_need, need)

            # Cap by available slide space
            per_tier_avail = (available - tier_gap * (num_tiers - 1)) / max(
                num_tiers, 1
            )
            tier_h = min(max_tier_need, per_tier_avail)
            tier_h = max(tier_h, 0.55)

            # Second pass: render
            for i, (label, title, combined) in enumerate(tier_items):
                tier_top = current_top + i * (tier_h + tier_gap)

                # Background card
                add_filled_box(slide, CONTENT_LEFT, tier_top, CONTENT_WIDTH, tier_h)
                # Label (e.g. "TIER 1: Discovery")
                label_h = _estimate_text_height(label.upper(), 10, inner_w, bold=True)
                add_text_box(
                    slide,
                    label.upper(),
                    left=CONTENT_LEFT + 0.15,
                    top=tier_top + 0.05,
                    width=inner_w,
                    height=label_h,
                    font_size=10,
                    bold=True,
                    color=self.accent_color,
                )
                # Title
                title_h = _estimate_text_height(title, 12, inner_w, bold=True)
                title_top = tier_top + 0.05 + label_h
                add_text_box(
                    slide,
                    title,
                    left=CONTENT_LEFT + 0.15,
                    top=title_top,
                    width=inner_w,
                    height=title_h,
                    font_size=12,
                    bold=True,
                    color=WHITE,
                )
                # Description + tokens
                desc_top = title_top + title_h + 0.02
                desc_h = tier_h - (desc_top - tier_top) - 0.05
                # Ensure desc box is tall enough for content
                if combined:
                    desc_est = _estimate_text_height(combined, 9, inner_w)
                    desc_h = max(desc_h, desc_est)
                if combined and desc_h > 0.08:
                    add_text_box(
                        slide,
                        combined,
                        left=CONTENT_LEFT + 0.15,
                        top=desc_top,
                        width=inner_w,
                        height=desc_h,
                        font_size=9,
                        color=GRAY_70,
                    )
                handled_elements.add(id(tiers[i]))
            current_top += num_tiers * (tier_h + tier_gap) + GAP_SECTION
            handled_elements.add(id(tier_stack))

        # ── Tier rows (cost-optimization style: name, uses, cost) ────────────
        tier_rows = [
            tr
            for tr in slide_div.find_all(class_="tier-row")
            if id(tr) not in handled_elements
        ]
        if tier_rows:
            # Header row
            cols_w = [1.8, 4.0, 2.6]
            tier_top = current_top
            for tr in tier_rows:
                name_el = tr.find(class_="tier-name")
                uses_el = tr.find(class_="tier-uses")
                cost_el = tr.find(class_="tier-cost")
                vals = [
                    get_text(name_el) if name_el else "",
                    get_text(uses_el) if uses_el else "",
                    get_text(cost_el) if cost_el else "",
                ]
                # Compute max row height across columns
                row_h = 0.35
                for j, val in enumerate(vals):
                    w = cols_w[j] if j < len(cols_w) else 2.0
                    h = _estimate_text_height(val, 12, w, bold=(j == 0))
                    row_h = max(row_h, h)
                left = CONTENT_LEFT
                for j, val in enumerate(vals):
                    w = cols_w[j] if j < len(cols_w) else 2.0
                    add_text_box(
                        slide,
                        val,
                        left=left,
                        top=tier_top,
                        width=w,
                        height=row_h,
                        font_size=12,
                        bold=(j == 0),
                        color=self.accent_color if j == 0 else GRAY_70,
                    )
                    left += w
                handled_elements.add(id(tr))
                tier_top += row_h + 0.05
            current_top = tier_top + GAP_SECTION

        # ── Diagram boxes (shadow-environments style flow) ───────────────────
        diagram = slide_div.find(class_="diagram")
        if diagram and id(diagram) not in handled_elements:
            boxes = diagram.find_all(class_="diagram-box")
            if boxes:
                num_boxes = len(boxes)
                gap_d = 0.15
                arrow_w = 0.45
                total_arrows = num_boxes - 1
                total_arrow_space = (
                    total_arrows * (arrow_w + 2 * gap_d) if total_arrows > 0 else 0
                )
                box_w = min(
                    2.4, (CONTENT_WIDTH - total_arrow_space) / max(num_boxes, 1)
                )
                total_w = num_boxes * box_w + total_arrow_space
                cur_left = CONTENT_LEFT + (CONTENT_WIDTH - total_w) / 2

                # Pre-scan boxes to compute uniform height
                inner_w = box_w - 0.16
                box_h = 0.80
                for box_el in boxes:
                    te = box_el.find(class_="diagram-box-title")
                    ce = box_el.find(class_="diagram-box-content")
                    th = (
                        _estimate_text_height(get_text(te), 12, inner_w, bold=True)
                        if te
                        else 0.30
                    )
                    ch = _estimate_text_height(get_text(ce), 10, inner_w) if ce else 0.0
                    box_h = max(box_h, th + ch + 0.24)

                for bi, box_el in enumerate(boxes):
                    title_el = box_el.find(class_="diagram-box-title")
                    content_el = box_el.find(class_="diagram-box-content")
                    add_filled_box(
                        slide,
                        cur_left,
                        current_top,
                        box_w,
                        box_h,
                        fill_color=DARK_GRAY,
                        border_color=self.accent_color,
                        border_width=1,
                    )
                    t_h = 0.30
                    if title_el:
                        t_text = get_text(title_el)
                        t_h = _estimate_text_height(t_text, 12, inner_w, bold=True)
                        t_h = max(0.30, t_h)
                        add_text_box(
                            slide,
                            t_text,
                            left=cur_left + 0.08,
                            top=current_top + 0.08,
                            width=inner_w,
                            height=t_h,
                            font_size=12,
                            bold=True,
                            color=WHITE,
                            align=PP_ALIGN.CENTER,
                        )
                    if content_el:
                        c_text = get_text(content_el)
                        c_h = _estimate_text_height(c_text, 10, inner_w)
                        c_h = max(0.25, min(c_h, box_h - t_h - 0.16))
                        add_text_box(
                            slide,
                            c_text,
                            left=cur_left + 0.08,
                            top=current_top + 0.08 + t_h,
                            width=inner_w,
                            height=c_h,
                            font_size=10,
                            color=GRAY_70,
                            align=PP_ALIGN.CENTER,
                        )
                    cur_left += box_w
                    if bi < num_boxes - 1:
                        add_text_box(
                            slide,
                            "\u2192",
                            left=cur_left + gap_d,
                            top=current_top + box_h / 2 - 0.18,
                            width=arrow_w,
                            height=0.40,
                            font_size=18,
                            bold=True,
                            color=self.accent_color,
                            align=PP_ALIGN.CENTER,
                        )
                        cur_left += arrow_w + 2 * gap_d
                    handled_elements.add(id(box_el))

                current_top += box_h + GAP_SECTION
            handled_elements.add(id(diagram))

        # ── Before/After comparison ──────────────────────────────────────────
        before_after = slide_div.find(class_="before-after")
        if before_after and id(before_after) not in handled_elements:
            ba_cards = before_after.find_all(
                class_=lambda c: c and ("before-card" in c or "after-card" in c)
            )
            col_w = CONTENT_WIDTH / 2 - 0.1
            for ci, ba_card in enumerate(ba_cards):
                ba_left = CONTENT_LEFT + ci * (col_w + 0.2)
                label_el = ba_card.find(class_="comparison-label")
                value_el = ba_card.find(class_="comparison-value")
                label = get_text(label_el) if label_el else ""
                value = get_text(value_el) if value_el else ""
                desc = get_text(ba_card)

                card_classes = ba_card.get("class", [])
                border_c = MS_ORANGE if "before-card" in card_classes else MS_GREEN

                add_filled_box(
                    slide,
                    ba_left,
                    current_top,
                    col_w,
                    1.2,
                    fill_color=DARK_GRAY,
                    border_color=border_c,
                    border_width=1,
                )
                label_h = _estimate_text_height(label, 14, col_w - 0.24, bold=True)
                label_h = max(0.30, label_h)
                add_text_box(
                    slide,
                    label,
                    left=ba_left + 0.12,
                    top=current_top + 0.08,
                    width=col_w - 0.24,
                    height=label_h,
                    font_size=14,
                    bold=True,
                    color=border_c,
                )
                if value:
                    val_h = _estimate_text_height(value, 28, col_w - 0.24, bold=True)
                    val_h = max(0.50, val_h)  # minimum for 28pt bold stats
                    add_text_box(
                        slide,
                        value,
                        left=ba_left + 0.12,
                        top=current_top + 0.38,
                        width=col_w - 0.24,
                        height=val_h,
                        font_size=28,
                        bold=True,
                        color=WHITE,
                    )
                # Remaining description below value
                remaining_text = desc.replace(label, "").replace(value, "").strip()
                if remaining_text and len(remaining_text) > 5:
                    desc_h = _estimate_text_height(remaining_text, 10, col_w - 0.24)
                    desc_h = max(0.30, min(desc_h, 0.45))
                    add_text_box(
                        slide,
                        remaining_text,
                        left=ba_left + 0.12,
                        top=current_top + 0.75,
                        width=col_w - 0.24,
                        height=desc_h,
                        font_size=10,
                        color=GRAY_70,
                    )
                handled_elements.add(id(ba_card))
            current_top += 1.3 + GAP_SECTION
            handled_elements.add(id(before_after))

        # ── Token display rows ───────────────────────────────────────────────
        token_displays = [
            td
            for td in slide_div.find_all(class_="token-display")
            if id(td) not in handled_elements
        ]
        if token_displays:
            row_h = 0.25
            for i, td in enumerate(token_displays):
                text = get_text(td)
                add_text_box(
                    slide,
                    text,
                    left=CONTENT_LEFT,
                    top=current_top + i * row_h,
                    width=CONTENT_WIDTH,
                    height=row_h,
                    font_size=11,
                    font_name=CODE_FONT,
                    color=WHITE,
                )
                handled_elements.add(id(td))
            current_top += len(token_displays) * row_h + GAP_NORMAL

        # ── Good/bad pattern lists ───────────────────────────────────────────
        good_patterns = [
            gp
            for gp in slide_div.find_all(class_="good-pattern")
            if id(gp) not in handled_elements
        ]
        bad_patterns = [
            bp
            for bp in slide_div.find_all(class_="bad-pattern")
            if id(bp) not in handled_elements
        ]
        if good_patterns or bad_patterns:
            col_w = CONTENT_WIDTH / 2 - 0.1
            # Bad patterns on left
            bad_top = current_top
            for bp in bad_patterns:
                bp_text = f"\u2717 {get_text(bp)}"
                bp_h = _estimate_text_height(bp_text, 13, col_w)
                bp_h = max(0.30, bp_h)
                add_text_box(
                    slide,
                    bp_text,
                    left=CONTENT_LEFT,
                    top=bad_top,
                    width=col_w,
                    height=bp_h,
                    font_size=13,
                    color=MS_RED,
                )
                handled_elements.add(id(bp))
                bad_top += bp_h + 0.02
            # Good patterns on right
            good_top = current_top
            for gp in good_patterns:
                gp_text = f"\u2713 {get_text(gp)}"
                gp_h = _estimate_text_height(gp_text, 13, col_w)
                gp_h = max(0.30, gp_h)
                add_text_box(
                    slide,
                    gp_text,
                    left=CONTENT_LEFT + col_w + 0.2,
                    top=good_top,
                    width=col_w,
                    height=gp_h,
                    font_size=13,
                    color=MS_GREEN,
                )
                handled_elements.add(id(gp))
                good_top += gp_h + 0.02
            current_top = max(bad_top, good_top) + GAP_SECTION

        # ── Summary rows (table-like summary) ────────────────────────────────
        summary_rows = [
            sr
            for sr in slide_div.find_all(class_="summary-row")
            if id(sr) not in handled_elements
        ]
        if summary_rows:
            sr_top = current_top
            for i, sr in enumerate(summary_rows):
                cells = sr.find_all(class_="summary-cell")
                num_cells = len(cells) or 1
                cell_w = CONTENT_WIDTH / num_cells
                is_first_row = i == 0
                fs = 12 if is_first_row else 11
                row_h = 0.30
                for j, cell in enumerate(cells):
                    text = get_text(cell)
                    cell_h = _estimate_text_height(text, fs, cell_w, bold=is_first_row)
                    row_h = max(row_h, cell_h)
                for j, cell in enumerate(cells):
                    text = get_text(cell)
                    add_text_box(
                        slide,
                        text,
                        left=CONTENT_LEFT + j * cell_w,
                        top=sr_top,
                        width=cell_w,
                        height=row_h,
                        font_size=fs,
                        bold=is_first_row,
                        color=self.accent_color if is_first_row else GRAY_70,
                    )
                    handled_elements.add(id(cell))
                handled_elements.add(id(sr))
                sr_top += row_h + 0.02
            current_top = sr_top + GAP_SECTION

        # ── Body text paragraphs ─────────────────────────────────────────────
        body_texts = [
            bt
            for bt in slide_div.find_all(class_="body-text")
            if id(bt) not in handled_elements
        ]
        for bt in body_texts:
            text = get_text(bt)
            if text:
                bt_h = _estimate_text_height(text, 14, CONTENT_WIDTH, bold=False)
                bt_h = max(0.3, min(bt_h, 3.5))
                add_text_box(
                    slide,
                    text,
                    left=CONTENT_LEFT,
                    top=current_top,
                    width=CONTENT_WIDTH,
                    height=bt_h,
                    font_size=14,
                    color=GRAY_70,
                    align=PP_ALIGN.CENTER,
                )
                current_top += bt_h + GAP_NORMAL
            handled_elements.add(id(bt))

        # ── Title meta (small info text on title slides) ─────────────────────
        title_meta = slide_div.find(class_="title-meta")
        if title_meta and id(title_meta) not in handled_elements:
            meta_text = get_text(title_meta)
            if meta_text:
                meta_top = max(current_top, 4.5)
                meta_h = _estimate_text_height(meta_text, 14, CONTENT_WIDTH, bold=False)
                meta_h = max(0.3, min(meta_h, 1.5))
                add_text_box(
                    slide,
                    meta_text,
                    left=CONTENT_LEFT,
                    top=meta_top,
                    width=CONTENT_WIDTH,
                    height=meta_h,
                    font_size=14,
                    color=GRAY_50,
                    align=PP_ALIGN.CENTER if is_centered else PP_ALIGN.LEFT,
                )
                current_top = meta_top + meta_h
            handled_elements.add(id(title_meta))

        # ── Highlight boxes ──────────────────────────────────────────────────
        highlight_boxes = slide_div.find_all(class_="highlight-box")
        for hb in highlight_boxes:
            if id(hb) in handled_elements:
                continue
            classes = hb.get("class", [])
            color = parse_color_from_class(classes) or self.accent_color
            rich = get_rich_text(hb)
            plain = get_text(hb)
            # Place at current_top; cap near bottom but never above current_top
            hb_top = min(current_top, SLIDE_HEIGHT - 0.85)
            add_highlight_box(
                slide,
                plain,
                top=hb_top,
                color=color,
                rich_runs=rich if len(rich) > 1 else None,
            )
            current_top = hb_top + 0.8
            handled_elements.add(id(hb))

        # ── Quote ────────────────────────────────────────────────────────────
        quote = slide_div.find(class_="quote")
        if quote:
            self._add_quote(slide, quote, current_top)
            current_top += 1.5  # quote height (1.2) + attribution (0.3)
            handled_elements.add(id(quote))

        # ── Small text (footer) ──────────────────────────────────────────────
        small_texts = slide_div.find_all(class_="small-text")
        small_text_top = max(current_top, 4.8)
        for st in small_texts:
            if id(st) in handled_elements:
                continue
            st_text = get_text(st)
            st_h = _estimate_text_height(st_text, 14, CONTENT_WIDTH, bold=False)
            st_h = max(0.3, min(st_h, 1.5))
            add_text_box(
                slide,
                st_text,
                left=CONTENT_LEFT,
                top=small_text_top,
                width=CONTENT_WIDTH,
                height=st_h,
                font_size=14,
                color=GRAY_50,
                align=PP_ALIGN.CENTER if is_centered else PP_ALIGN.LEFT,
            )
            small_text_top += st_h
            handled_elements.add(id(st))
        current_top = small_text_top

        # ── Fallback: render unrecognized elements with text content ─────────
        # Instead of silently dropping unknown elements, render them as
        # generic text boxes so no content vanishes.
        # Build ancestor skip set: any element that is an ancestor of a
        # handled element should also be skipped to prevent duplication.
        _ancestor_skip: set[int] = set()
        for _desc in slide_div.descendants:
            if isinstance(_desc, Tag) and id(_desc) in handled_elements:
                _par = _desc.parent
                while _par and _par is not slide_div:
                    _ancestor_skip.add(id(_par))
                    _par = _par.parent
        for el in slide_div.children:
            if not isinstance(el, Tag):
                continue
            if id(el) in handled_elements or id(el) in _ancestor_skip:
                continue
            # Skip elements that are parents of already-handled children
            if any(
                id(desc) in handled_elements
                for desc in el.descendants
                if isinstance(desc, Tag)
            ):
                continue
            if hasattr(el, "get_text") and el.get_text(strip=True):
                text = get_text(el)
                if text and len(text) > 5:  # skip trivial fragments
                    text = _truncate_to_fit(text, 14, CONTENT_WIDTH, 5.5)
                    fb_height = _estimate_text_height(
                        text, 14, CONTENT_WIDTH, bold=False
                    )
                    fb_height = max(0.3, fb_height)
                    add_text_box(
                        slide,
                        text,
                        left=CONTENT_LEFT,
                        top=current_top,
                        width=CONTENT_WIDTH,
                        height=fb_height,
                        font_size=14,
                        color=GRAY_70,
                    )
                    current_top += fb_height + GAP_NORMAL
                    handled_elements.add(id(el))

        # ── Overflow compression ─────────────────────────────────────────────
        # When content extends past slide height, proportionally compress
        # shape positions and heights to fit within the usable area.
        # Check both cursor AND actual shape bounds (shapes may extend
        # past current_top when given generous content-based heights).
        actual_bottom = current_top
        for shape in slide.shapes:
            s_bot = (shape.top + shape.height) / 914400
            actual_bottom = max(actual_bottom, s_bot)
        if actual_bottom > SLIDE_HEIGHT:
            usable = SLIDE_HEIGHT - 0.10  # leave 0.1" bottom margin
            # Find content bounds (skip shapes at very top like backgrounds)
            shape_tops = []
            shape_bots = []
            for shape in slide.shapes:
                s_top = shape.top / 914400
                s_bot = (shape.top + shape.height) / 914400
                if s_top >= 0.3:  # only content shapes, not backgrounds
                    shape_tops.append(s_top)
                    shape_bots.append(s_bot)

            if shape_tops and shape_bots:
                content_start = min(shape_tops)
                content_end = max(shape_bots)
                content_height = content_end - content_start

                if content_height > 0:
                    scale = (usable - content_start) / content_height

                    if 0.40 < scale < 1.0:
                        for shape in slide.shapes:
                            s_top = shape.top / 914400
                            if s_top >= content_start:
                                rel_top = s_top - content_start
                                new_top = content_start + rel_top * scale
                                shape.top = int(new_top * 914400)
                                # Do NOT compress heights — that causes text
                                # overflow inside boxes.  Position-only
                                # compression closes gaps without shrinking
                                # text containers.
                    elif scale <= 0.40:
                        # Extreme overflow — position-only compression at
                        # 40% floor (don't scale heights to avoid creating
                        # text overflow inside boxes).
                        floor_scale = 0.40
                        for shape in slide.shapes:
                            s_top = shape.top / 914400
                            if s_top >= content_start:
                                rel_top = s_top - content_start
                                new_top = content_start + rel_top * floor_scale
                                shape.top = int(new_top * 914400)
                        self.warnings.append(
                            f"Slide {slide_num}: severe overflow "
                            f'({current_top:.1f}"), compressed at 40% floor '
                            f"(needed {scale:.0%} scale)"
                        )

        # ── Auto-shrink safety net ──────────────────────────────────────

    # ── Card layouts ───────────────────────────────────────────────────── ─────────────────────────────────────────────────────────

    def _add_cards(
        self, slide, cards: list[Tag], top: float, container: Optional[Tag] = None
    ) -> float:
        """Add a row of cards to the slide. Returns total height consumed."""
        num_cards = len(cards)
        if num_cards == 0:
            return 0.0

        # Determine column count from container class.
        #
        # CSS auto-fit classes ("halves", "thirds") use:
        #   grid-template-columns: repeat(auto-fit, minmax(min(Npx, 100%), 1fr))
        # At presentation viewport (~1120px content), "halves" (min 300px) fits
        # up to 3 columns; "thirds" (min 280px) fits up to 3.  Explicit grid-N
        # classes always force exactly N columns.  For auto-fit classes we cap at
        # the CSS-implied maximum and let card count fill rows naturally.
        #
        # Multi-row logic: when >4 cards, cap columns so each card is ≥2.5" wide.
        # E.g. grid-5 with 5 cards → 3+2 layout; grid with 6 → 3+3.
        forced_cols = None
        if container:
            cls = " ".join(container.get("class", []))
            if "grid-5" in cls or "fifths" in cls:
                # 5 cols at 8.4" width → 1.52" each (too narrow).
                # Cap at 3 cols for ≥2.5" minimum card width.
                forced_cols = 3
            elif "grid-4" in cls or "fourths" in cls:
                forced_cols = 4
            elif "grid-3" in cls or "thirds" in cls:
                forced_cols = 3
            elif "grid-2" in cls:
                forced_cols = 2
            elif "halves" in cls:
                # auto-fit with minmax(300px, 1fr) → up to 3 cols at 1120px
                forced_cols = min(num_cards, 3)

        cols = forced_cols or min(num_cards, 4)

        # Enforce minimum card width of 2.5" — reduce columns if needed
        gap = 0.2
        min_card_width = 2.5
        while cols > 1 and (CONTENT_WIDTH - gap * (cols - 1)) / cols < min_card_width:
            cols -= 1

        num_rows = -(-num_cards // cols)  # ceil division
        gap = 0.2
        total_width = CONTENT_WIDTH
        card_width = (total_width - gap * (cols - 1)) / cols
        start_left = CONTENT_LEFT

        # Adaptive card height: estimate from content, fit within slide space
        available_height = SLIDE_HEIGHT - top - 0.3  # 0.3" bottom margin
        row_gap = 0.2  # gap between rows of cards

        # Compute needed height for each card based on actual content
        inner_w = card_width - 0.3  # text area inside card (0.15" pad each side)
        max_needed = 0.8  # minimum card height
        for card_el in cards:
            card_classes = card_el.get("class", [])
            if "module-card" in card_classes:
                # Module cards have fixed proportional layout — use default
                max_needed = max(max_needed, 1.8)
                continue
            title_el = card_el.find(class_=["card-title", "tool-name"])
            text_el = card_el.find(class_=["card-text", "card-desc", "tool-desc"])
            usage_el = card_el.find(class_="tool-usage")
            number_el = card_el.find(class_="card-number")
            t_title = get_text(title_el) if title_el else ""
            t_text = get_text(text_el) if text_el else ""
            if usage_el:
                u = get_text(usage_el)
                if u:
                    t_text = f"{t_text}\n{u}" if t_text else u
            if number_el:
                # number cards need big number + title + desc
                num_text = get_text(number_el)
                num_font = max(24, min(48, int(1.8 * 28)))
                text_est = _estimate_text_height(t_text, 10, inner_w) if t_text else 0.0
                need = (
                    0.08  # top pad
                    + _estimate_text_height(num_text, num_font, inner_w, bold=True)
                    + _estimate_text_height(t_title, 14, inner_w, bold=True)
                    + text_est
                    + 0.08  # bottom pad
                )
            else:
                title_h = _estimate_text_height(t_title, 16, inner_w, bold=True) + 0.05
                text_h = _estimate_text_height(t_text, 12, inner_w) if t_text else 0.0
                need = 0.15 + title_h + GAP_TIGHT + text_h + 0.05
            max_needed = max(max_needed, need)

        if num_rows == 1:
            card_height = min(max(max_needed, 0.8), available_height)
        else:
            per_row_avail = (available_height - row_gap * (num_rows - 1)) / num_rows
            card_height = min(max(max_needed, 0.8), per_row_avail)
        card_height = max(card_height, 0.8)  # never smaller than 0.8"

        # Center cards if fewer than cols
        if num_cards < cols:
            actual_width = num_cards * card_width + (num_cards - 1) * gap
            start_left = CONTENT_LEFT + (total_width - actual_width) / 2

        for i, card_el in enumerate(cards):
            col = i % cols
            row = i // cols
            left = start_left + col * (card_width + gap)
            row_top = top + row * (card_height + row_gap)

            # Detect card type
            card_classes = card_el.get("class", [])
            is_module_card = "module-card" in card_classes
            if is_module_card:
                self._add_module_card(
                    slide, card_el, left, row_top, card_width, card_height
                )
            else:
                title_el = card_el.find(class_=["card-title", "tool-name"])
                text_el = card_el.find(class_=["card-text", "card-desc", "tool-desc"])
                usage_el = card_el.find(class_="tool-usage")
                number_el = card_el.find(class_="card-number")

                title = get_text(title_el) if title_el else ""
                text = get_text(text_el) if text_el else ""
                # Append tool-usage text (e.g. example invocations)
                if usage_el:
                    usage_text = get_text(usage_el)
                    if usage_text:
                        text = f"{text}\n{usage_text}" if text else usage_text
                rich = get_rich_text(text_el) if text_el else []

                if number_el:
                    number = get_text(number_el)
                    self._add_number_card(
                        slide,
                        number,
                        title,
                        text,
                        left,
                        row_top,
                        card_width,
                        card_height,
                    )
                else:
                    add_card(
                        slide,
                        title,
                        text,
                        left,
                        row_top,
                        width=card_width,
                        height=card_height,
                        title_color=self.accent_color,
                        rich_runs=rich if len(rich) > 1 else None,
                    )

        # Return total height consumed by all rows
        return num_rows * card_height + (num_rows - 1) * row_gap

    def _add_principles(self, slide, principles: list[Tag], top: float) -> float:
        """Add numbered principle boxes in a two-column tenet layout.

        Returns the row_height used (content-adaptive).
        """
        cols = 2
        col_width = CONTENT_WIDTH / 2 - 0.1
        text_inner_w = col_width - 0.35  # text area width inside tenet
        row_gap = 0.1

        # First pass: extract content and compute needed row height
        items: list[tuple[str, str]] = []
        max_needed = 0.85  # minimum row height
        for i, principle in enumerate(principles):
            num_el = principle.find(class_=["principle-number", "principle-num"])
            content_el = principle.find(class_=["principle-content", "principle-text"])
            number = get_text(num_el) if num_el else str(i + 1)
            title = ""
            desc = ""
            if content_el:
                h3 = content_el.find("h3")
                p_tag = content_el.find("p")
                strong = content_el.find("strong")
                if h3:
                    title = get_text(h3)
                    desc = get_text(p_tag) if p_tag else ""
                elif strong:
                    # Pattern: <strong>Title</strong> Description text
                    title = get_text(strong)
                    strong.extract()
                    desc = get_text(content_el).strip()
                elif p_tag:
                    title = get_text(p_tag)
                else:
                    title = get_text(content_el)

            full_title = f"{number}. {title}"
            items.append((full_title, desc))
            # Estimate: top pad + title + gap + description + bottom pad
            title_h = _estimate_text_height(full_title, 14, text_inner_w, bold=True)
            desc_h = _estimate_text_height(desc, 11, text_inner_w) if desc else 0.0
            need = 0.08 + title_h + GAP_TIGHT + desc_h + 0.05
            max_needed = max(max_needed, need)

        row_height = max_needed

        # Second pass: place the tenet boxes
        for i, (full_title, desc) in enumerate(items):
            col = i % cols
            row = i // cols
            left = CONTENT_LEFT + col * (col_width + 0.2)
            ptop = top + row * (row_height + row_gap)

            add_tenet(
                slide,
                full_title,
                desc,
                left,
                ptop,
                width=col_width,
                height=row_height,
                accent_color=self.accent_color,
            )

        return row_height

    def _add_module_card(
        self,
        slide,
        card_el: Tag,
        left: float,
        top: float,
        width: float,
        height: float = 1.8,
    ):
        """Add a module-card (name, contract, purpose) with accent top border."""

        # Card background with top accent border
        add_filled_box(slide, left, top, width, height)
        # Top accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.accent_color
        bar.line.fill.background()

        # Scale internal layout using content-aware heights
        pad = 0.10
        inner_w = width - 0.24
        y_cursor = top + pad

        # Module name
        name_el = card_el.find(class_="module-name")
        name_fs = max(11, min(15, int(height * 9)))
        if name_el:
            name_text = get_text(name_el)
            name_h = _estimate_text_height(name_text, name_fs, inner_w, bold=True)
            name_h = max(0.25, min(name_h, height * 0.30))
            add_text_box(
                slide,
                name_text,
                left=left + 0.12,
                top=y_cursor,
                width=inner_w,
                height=name_h,
                font_size=name_fs,
                bold=True,
                color=self.accent_color,
            )
            y_cursor += name_h

        # Contract (monospace)
        contract_el = card_el.find(class_="module-contract")
        contract_fs = max(8, min(10, int(height * 6)))
        if contract_el:
            contract_text = get_text(contract_el)
            contract_h = _estimate_text_height(contract_text, contract_fs, inner_w)
            contract_h = max(0.20, min(contract_h, height * 0.25))
            add_text_box(
                slide,
                contract_text,
                left=left + 0.12,
                top=y_cursor,
                width=inner_w,
                height=contract_h,
                font_size=contract_fs,
                font_name=CODE_FONT,
                color=CODE_GREEN,
            )
            y_cursor += contract_h

        # Purpose
        purpose_el = card_el.find(class_="module-purpose")
        remaining = height - (y_cursor - top) - pad
        if purpose_el and remaining > 0.1:
            add_text_box(
                slide,
                get_text(purpose_el),
                left=left + 0.12,
                top=y_cursor,
                width=width - 0.24,
                height=remaining,
                font_size=max(8, min(11, int(height * 6))),
                color=GRAY_70,
            )

    def _add_number_card(
        self,
        slide,
        number: str,
        title: str,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float = 1.8,
    ):
        """Add a card with a big number."""
        add_filled_box(slide, left, top, width, height)

        # Scale internal layout proportionally to card height
        pad = 0.08
        inner_w = width - 0.2
        num_font = max(24, min(48, int(height * 28)))
        title_font = max(10, min(14, int(height * 8)))
        text_font = max(8, min(10, int(height * 6)))

        # Estimate needed heights from actual content
        num_h = _estimate_text_height(number, num_font, inner_w, bold=True)
        num_h = max(0.5, min(num_h, height * 0.50))
        title_h = _estimate_text_height(title, title_font, inner_w, bold=True)
        title_h = max(0.25, min(title_h, height * 0.25))

        add_text_box(
            slide,
            number,
            left=left + 0.1,
            top=top + pad,
            width=inner_w,
            height=num_h,
            font_size=num_font,
            bold=True,
            color=MS_CYAN,
            align=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            title,
            left=left + 0.1,
            top=top + pad + num_h,
            width=inner_w,
            height=title_h,
            font_size=title_font,
            bold=True,
            color=self.accent_color,
            align=PP_ALIGN.CENTER,
        )
        remaining = height - pad - num_h - title_h - pad
        if remaining > 0.1 and text:
            add_text_box(
                slide,
                text,
                left=left + 0.1,
                top=top + pad + num_h + title_h,
                width=inner_w,
                height=remaining,
                font_size=text_font,
                color=GRAY_70,
                align=PP_ALIGN.CENTER,
            )

    # ── Code blocks ──────────────────────────────────────────────────────────

    def _add_code_block(self, slide, code_el: Tag, top: float):
        """Add a code block with syntax highlighting."""
        # Extract structured code runs from HTML spans
        code_runs = self._extract_code_runs(code_el)
        plain_text = get_text(code_el)
        # Estimate height from actual content; monospace chars are ~20% wider
        # than Arial so multiply estimate by 1.15 to account for extra wrapping.
        text_h = _estimate_text_height(
            plain_text, 10, CONTENT_WIDTH - 0.4, line_spacing=1.15
        )
        height = max(1.0, min(text_h * 1.15 + 0.24, 4.5))

        # Dark background
        add_filled_box(
            slide,
            CONTENT_LEFT,
            top,
            CONTENT_WIDTH,
            height,
            fill_color=CODE_BG,
            border_color=RGBColor(0x30, 0x30, 0x30),
            border_width=1,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )

        # Code text with syntax highlighting
        if code_runs:
            box = slide.shapes.add_textbox(
                Inches(CONTENT_LEFT + 0.2),
                Inches(top + 0.12),
                Inches(CONTENT_WIDTH - 0.4),
                Inches(height - 0.24),
            )
            tf = box.text_frame
            tf.word_wrap = True

            # Split runs into lines (paragraphs)
            paragraphs: list[list[dict]] = [[]]
            for r in code_runs:
                parts = r["text"].split("\n")
                for part_idx, part in enumerate(parts):
                    if part_idx > 0:
                        paragraphs.append([])
                    if part:
                        paragraphs[-1].append({**r, "text": part})

            for p_idx, p_runs in enumerate(paragraphs):
                p = tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                if not p_runs:
                    run = p.add_run()
                    run.text = ""
                    _set_font(run, name=CODE_FONT, size=10, color=CODE_DEFAULT)
                else:
                    for r in p_runs:
                        run = p.add_run()
                        run.text = r["text"]
                        _set_font(
                            run,
                            name=CODE_FONT,
                            size=10,
                            bold=r.get("bold", False),
                            color=r.get("color", CODE_DEFAULT),
                        )
        else:
            # Fallback: plain monospace text
            add_text_box(
                slide,
                plain_text,
                left=CONTENT_LEFT + 0.2,
                top=top + 0.12,
                width=CONTENT_WIDTH - 0.4,
                height=height - 0.24,
                font_size=10,
                font_name=CODE_FONT,
                color=CODE_DEFAULT,
            )

    def _add_code_block_sized(
        self, slide, code_el: Tag, top: float, left: float, width: float
    ):
        """Add a code block at a specific position and width (for side-by-side grids)."""
        code_runs = self._extract_code_runs(code_el)
        plain_text = get_text(code_el)
        # Estimate height from content; 1.15× for monospace char-width difference
        pad = 0.15
        text_h = _estimate_text_height(
            plain_text, 9, width - 2 * pad, line_spacing=1.15
        )
        height = max(1.0, min(text_h * 1.15 + 0.20, 3.5))

        # Dark background
        add_filled_box(
            slide,
            left,
            top,
            width,
            height,
            fill_color=CODE_BG,
            border_color=RGBColor(0x30, 0x30, 0x30),
            border_width=1,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )

        pad = 0.15
        font_size = 9  # slightly smaller for side-by-side

        if code_runs:
            box = slide.shapes.add_textbox(
                Inches(left + pad),
                Inches(top + 0.10),
                Inches(width - 2 * pad),
                Inches(height - 0.20),
            )
            tf = box.text_frame
            tf.word_wrap = True

            paragraphs: list[list[dict]] = [[]]
            for r in code_runs:
                parts = r["text"].split("\n")
                for part_idx, part in enumerate(parts):
                    if part_idx > 0:
                        paragraphs.append([])
                    if part:
                        paragraphs[-1].append({**r, "text": part})

            for p_idx, p_runs in enumerate(paragraphs):
                p = tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                if not p_runs:
                    run = p.add_run()
                    run.text = ""
                    _set_font(run, name=CODE_FONT, size=font_size, color=CODE_DEFAULT)
                else:
                    for r in p_runs:
                        run = p.add_run()
                        run.text = r["text"]
                        _set_font(
                            run,
                            name=CODE_FONT,
                            size=font_size,
                            bold=r.get("bold", False),
                            color=r.get("color", CODE_DEFAULT),
                        )
        else:
            add_text_box(
                slide,
                plain_text,
                left=left + pad,
                top=top + 0.10,
                width=width - 2 * pad,
                height=height - 0.20,
                font_size=font_size,
                font_name=CODE_FONT,
                color=CODE_DEFAULT,
            )

    def _extract_code_runs(self, code_el: Tag) -> list[dict]:
        """Extract syntax-highlighted runs from a code block element."""
        el_copy = deepcopy(code_el)
        _replace_br_tags(el_copy)

        color_map = {
            "code-keyword": CODE_BLUE,
            "keyword": CODE_BLUE,
            "code-string": CODE_STRING,
            "string": CODE_STRING,
            "code-comment": CODE_GRAY,
            "comment": CODE_GRAY,
            "code-type": CODE_GREEN,
            "type": CODE_GREEN,
            "code-func": CODE_YELLOW,
            "func": CODE_YELLOW,
            "code-number": CODE_PURPLE,
            "number": CODE_PURPLE,
            "layer-kernel": CODE_BLUE,
            "layer-foundation": CODE_GREEN,
            "layer-apps": CODE_PURPLE,
            "layer-modules": CODE_YELLOW,
        }

        runs: list[dict] = []
        for child in el_copy.descendants:
            if isinstance(child, NavigableString):
                text = str(child)
                if not text:
                    continue
                # Determine color from parent span class
                color = CODE_DEFAULT
                bold = False
                parent = child.parent
                while parent and parent != el_copy:
                    if parent.name == "span":
                        for cls in parent.get("class", []):
                            if cls in color_map:
                                color = color_map[cls]
                                break
                    if parent.name in ("strong", "b"):
                        bold = True
                    parent = parent.parent

                runs.append({"text": text, "color": color, "bold": bold})

        # Merge adjacent runs with identical formatting
        merged: list[dict] = []
        for r in runs:
            if (
                merged
                and merged[-1]["color"] == r["color"]
                and merged[-1]["bold"] == r["bold"]
            ):
                merged[-1]["text"] += r["text"]
            else:
                merged.append(r)

        return merged

    # ── Flow diagrams ────────────────────────────────────────────────────────

    def _add_flow_step_content(
        self, slide, step: Tag, left: float, top: float, width: float, height: float
    ):
        """Render the interior of a single flow/workflow step box."""
        # Extended class lookups covering both flow-* and step-* conventions
        num_el = step.find(class_=["step-number"])
        title_el = step.find(
            class_=["flow-step-title", "workflow-step-title", "step-title"]
        )
        desc_el = step.find(
            class_=["flow-step-desc", "workflow-step-desc", "step-desc"]
        )
        turns_el = step.find(class_="step-turns")

        y = top + 0.08
        pad_x = 0.12
        inner_w = width - 2 * pad_x

        if num_el and title_el:
            # Number + Title on one line
            _ft_text = f"{get_text(num_el)}. {get_text(title_el)}"
            _ft_h = _estimate_text_height(_ft_text, 14, inner_w, bold=True)
            _ft_h = max(0.30, min(_ft_h, 0.60))
            add_text_box(
                slide,
                _ft_text,
                left=left + pad_x,
                top=y,
                width=inner_w,
                height=_ft_h,
                font_size=14,
                bold=True,
                color=WHITE,
            )
            y += _ft_h
        elif title_el:
            _ft_text = get_text(title_el)
            _ft_h = _estimate_text_height(_ft_text, 14, inner_w, bold=True)
            _ft_h = max(0.30, min(_ft_h, 0.60))
            add_text_box(
                slide,
                _ft_text,
                left=left + pad_x,
                top=y,
                width=inner_w,
                height=_ft_h,
                font_size=14,
                bold=True,
                color=WHITE,
                align=PP_ALIGN.CENTER,
            )
            y += _ft_h

        if desc_el:
            desc_text = get_text(desc_el)
            desc_needed = _estimate_text_height(desc_text, 11, inner_w)
            # Use estimated height — allow text box to extend if step box is tight
            desc_h = desc_needed
            if desc_h > 0.10:
                add_text_box(
                    slide,
                    desc_text,
                    left=left + pad_x,
                    top=y,
                    width=inner_w,
                    height=desc_h,
                    font_size=11,
                    color=GRAY_70,
                )
                y += desc_h

        if turns_el:
            turns_text = get_text(turns_el)
            turns_needed = _estimate_text_height(turns_text, 10, inner_w)
            remaining = height - (y - top) - 0.05
            turns_h = max(remaining, turns_needed)
            if turns_h > 0.10:
                add_text_box(
                    slide,
                    turns_text,
                    left=left + pad_x,
                    top=y,
                    width=inner_w,
                    height=turns_h,
                    font_size=10,
                    italic=True,
                    color=self.accent_color,
                )

        # Fallback: no structured children recognised
        if not title_el and not num_el:
            add_text_box(
                slide,
                get_text(step),
                left=left + 0.08,
                top=top + 0.08,
                width=width - 0.16,
                height=height - 0.16,
                font_size=11,
                color=WHITE,
                align=PP_ALIGN.CENTER,
            )

    def _add_flow_diagram(self, slide, flow_el: Tag, top: float) -> float:
        """Add a flow diagram. Returns total height consumed."""
        # Collect steps (flow-box, flow-step, workflow-step)
        steps = flow_el.find_all(class_=["flow-box", "flow-step", "workflow-step"])
        if not steps:
            return 0.0

        num_steps = len(steps)

        # ── Multi-row grid when >4 steps (keeps boxes ≥2.5" wide) ────────
        if num_steps > 4:
            cols = 3
            num_rows = -(-num_steps // cols)  # ceil division
            gap = 0.20
            row_gap = 0.25
            box_width = (CONTENT_WIDTH - gap * (cols - 1)) / cols

            # Content-aware box height: scan all steps for needed height
            pad_x = 0.12
            inner_w = box_width - 2 * pad_x
            max_step_need = 1.05  # minimum
            for step in steps:
                need = 0.08  # top pad
                t_el = step.find(
                    class_=["flow-step-title", "workflow-step-title", "step-title"]
                )
                n_el = step.find(class_="step-number")
                if t_el or n_el:
                    need += 0.30  # title row
                d_el = step.find(
                    class_=["flow-step-desc", "workflow-step-desc", "step-desc"]
                )
                if d_el:
                    need += _estimate_text_height(get_text(d_el), 11, inner_w)
                tu_el = step.find(class_="step-turns")
                if tu_el:
                    need += _estimate_text_height(get_text(tu_el), 10, inner_w)
                need += 0.05  # bottom pad
                max_step_need = max(max_step_need, need)
            # Cap by available slide space
            avail_per_row = (SLIDE_HEIGHT - top - 0.3 - (num_rows - 1) * row_gap) / max(
                num_rows, 1
            )
            box_height = min(max_step_need, avail_per_row)

            for i, step in enumerate(steps):
                col = i % cols
                row = i // cols
                left = CONTENT_LEFT + col * (box_width + gap)
                rtop = top + row * (box_height + row_gap)

                add_filled_box(
                    slide,
                    left,
                    rtop,
                    box_width,
                    box_height,
                    fill_color=DARK_GRAY,
                    border_color=self.accent_color,
                    border_width=1,
                )
                self._add_flow_step_content(
                    slide, step, left, rtop, box_width, box_height
                )

                # Horizontal arrow to next box in same row (not at row end)
                if col < cols - 1 and i < num_steps - 1:
                    arrow_x = left + box_width + gap / 2 - 0.22
                    add_text_box(
                        slide,
                        "\u2192",
                        left=arrow_x,
                        top=rtop + box_height / 2 - 0.22,
                        width=0.45,
                        height=0.45,
                        font_size=16,
                        bold=True,
                        color=self.accent_color,
                        align=PP_ALIGN.CENTER,
                    )

            return num_rows * box_height + (num_rows - 1) * row_gap

        # ── Single-row flow (≤4 steps) ───────────────────────────────────
        gap = 0.15
        arrow_width = 0.45
        total_arrows = num_steps - 1
        total_arrow_space = (
            total_arrows * (arrow_width + 2 * gap) if total_arrows > 0 else 0
        )
        remaining = CONTENT_WIDTH - total_arrow_space
        box_width = remaining / num_steps
        box_width = min(box_width, 2.5)

        # Content-aware height (same scan as multi-row path)
        pad_x = 0.12
        sr_inner_w = box_width - 2 * pad_x
        box_height = 0.9  # minimum
        for step in steps:
            need = 0.08
            t_el = step.find(
                class_=["flow-step-title", "workflow-step-title", "step-title"]
            )
            n_el = step.find(class_="step-number")
            if t_el or n_el:
                need += 0.30
            d_el = step.find(
                class_=["flow-step-desc", "workflow-step-desc", "step-desc"]
            )
            if d_el:
                need += _estimate_text_height(get_text(d_el), 11, sr_inner_w)
            tu_el = step.find(class_="step-turns")
            if tu_el:
                need += _estimate_text_height(get_text(tu_el), 10, sr_inner_w)
            need += 0.05
            box_height = max(box_height, need)
        # Cap to available slide space
        box_height = min(box_height, SLIDE_HEIGHT - top - 0.3)

        # Center the whole diagram
        total_width = num_steps * box_width + total_arrow_space
        start_left = CONTENT_LEFT + (CONTENT_WIDTH - total_width) / 2

        current_left = start_left
        for i, step in enumerate(steps):
            add_filled_box(
                slide,
                current_left,
                top,
                box_width,
                box_height,
                fill_color=DARK_GRAY,
                border_color=self.accent_color,
                border_width=1,
            )
            self._add_flow_step_content(
                slide, step, current_left, top, box_width, box_height
            )

            current_left += box_width

            # Arrow between steps
            if i < num_steps - 1:
                arrow_left = current_left + gap
                arrow_top = top + box_height / 2 - 0.22
                add_text_box(
                    slide,
                    "\u2192",
                    left=arrow_left,
                    top=arrow_top,
                    width=arrow_width,
                    height=0.45,
                    font_size=20,
                    bold=True,
                    color=self.accent_color,
                    align=PP_ALIGN.CENTER,
                )
                current_left += arrow_width + 2 * gap

        return box_height

    # ── Notification stacks ──────────────────────────────────────────────────

    def _add_notification_stack(self, slide, stack_el: Tag, top: float):
        """Add a notification stack (allowed/blocked notification items)."""
        notifications = stack_el.find_all(class_="notification")
        if not notifications:
            return

        row_height = 0.55
        row_gap = 0.08
        total_width = 6.0
        start_left = CONTENT_LEFT + (CONTENT_WIDTH - total_width) / 2

        for i, notif in enumerate(notifications):
            ntop = top + i * (row_height + row_gap)
            classes = notif.get("class", [])
            is_allowed = "allowed" in classes
            is_blocked = "blocked" in classes

            # Background
            if is_allowed:
                bg = RGBColor(0x0A, 0x1A, 0x0D)
                border = MS_GREEN
            elif is_blocked:
                bg = RGBColor(0x1A, 0x0A, 0x0A)
                border = MS_RED
            else:
                bg = DARK_GRAY
                border = BORDER_GRAY

            add_filled_box(
                slide,
                start_left,
                ntop,
                total_width,
                row_height,
                fill_color=bg,
                border_color=border,
                border_width=1,
                shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
            )

            # Status icon
            icon = "\u2713" if is_allowed else "\u2715"
            icon_color = MS_GREEN if is_allowed else MS_RED
            add_text_box(
                slide,
                icon,
                left=start_left + 0.1,
                top=ntop + 0.05,
                width=0.35,
                height=0.45,
                font_size=16,
                bold=True,
                color=icon_color,
                align=PP_ALIGN.CENTER,
            )

            # Notification title + body
            title_el = notif.find(class_="notification-title")
            body_el = notif.find(class_="notification-body")
            title = get_text(title_el) if title_el else ""
            body = get_text(body_el) if body_el else ""

            notif_w = total_width - 0.70
            title_h = _estimate_text_height(title, 12, notif_w, bold=True)
            title_h = max(0.22, min(title_h, row_height * 0.5))
            add_text_box(
                slide,
                title,
                left=start_left + 0.5,
                top=ntop + 0.04,
                width=notif_w,
                height=title_h,
                font_size=12,
                bold=True,
                color=WHITE,
            )
            body_h = _estimate_text_height(body, 10, notif_w)
            body_h = max(0.20, min(body_h, row_height - title_h - 0.08))
            add_text_box(
                slide,
                _truncate_to_fit(body, 10, notif_w, body_h),
                left=start_left + 0.5,
                top=ntop + 0.04 + title_h,
                width=notif_w,
                height=body_h,
                font_size=10,
                color=GRAY_70,
            )

    # ── Comparison table (two-column grid, not HTML <table>) ─────────────────

    def _add_comparison_table(self, slide, comp_el: Tag, top: float):
        """Add a comparison-table (CSS grid with .header .left .right cells)."""
        children = [c for c in comp_el.children if isinstance(c, Tag)]
        if not children:
            return

        col_width = CONTENT_WIDTH / 2
        current_row_top = top

        # Pre-scan to pair left/right cells and compute row heights
        # Children alternate: left, right, left, right...
        pending_left_text: Optional[str] = None
        left_h: float = 0.32  # default; updated when left cell is seen

        for child in children:
            classes = child.get("class", [])
            text = get_text(child)
            is_header = "header" in classes
            is_left = "left" in classes
            is_right = "right" in classes

            fs = 13 if is_header else 11
            cell_h = _estimate_text_height(text, fs, col_width, bold=is_header)
            cell_h = max(0.28, min(cell_h, 0.8))

            if is_left:
                # Compute left-side height; store for pairing with right
                left_h = cell_h
                pending_left_text = text

            if is_header and is_left:
                add_text_box(
                    slide,
                    text,
                    left=CONTENT_LEFT,
                    top=current_row_top,
                    width=col_width,
                    height=cell_h,
                    font_size=13,
                    bold=True,
                    color=self.accent_color,
                )
            elif is_header and is_right:
                # Use max of left and right heights for this row
                row_h = max(left_h, cell_h) if pending_left_text is not None else cell_h
                add_text_box(
                    slide,
                    text,
                    left=CONTENT_LEFT + col_width,
                    top=current_row_top,
                    width=col_width,
                    height=row_h,
                    font_size=13,
                    bold=True,
                    color=CODE_GREEN,
                )
                current_row_top += row_h + 0.05
                pending_left_text = None
            elif is_left and not is_header:
                add_text_box(
                    slide,
                    text,
                    left=CONTENT_LEFT,
                    top=current_row_top,
                    width=col_width,
                    height=cell_h,
                    font_size=11,
                    color=WHITE,
                )
            elif is_right and not is_header:
                row_h = max(left_h, cell_h) if pending_left_text is not None else cell_h
                add_text_box(
                    slide,
                    text,
                    left=CONTENT_LEFT + col_width,
                    top=current_row_top,
                    width=col_width,
                    height=row_h,
                    font_size=11,
                    color=GRAY_70,
                )
                current_row_top += row_h
                pending_left_text = None

    # ── Architecture diagram ─────────────────────────────────────────────────

    def _add_architecture_diagram(self, slide, arch_el: Tag, top: float):
        """Add an architecture diagram (text art with colored spans)."""
        runs = self._extract_code_runs(arch_el)
        if not runs:
            plain = get_text(arch_el)
            add_text_box(
                slide,
                plain,
                left=CONTENT_LEFT,
                top=top,
                width=CONTENT_WIDTH,
                height=3.0,
                font_size=11,
                font_name=CODE_FONT,
                color=CODE_DEFAULT,
            )
            return

        # Use code block rendering
        plain_text = get_text(arch_el)
        lines = plain_text.count("\n") + 1
        height = max(1.5, min(lines * 0.18 + 0.3, 3.5))

        add_filled_box(
            slide,
            CONTENT_LEFT,
            top,
            CONTENT_WIDTH,
            height,
            fill_color=CODE_BG,
            border_color=RGBColor(0x30, 0x30, 0x30),
            border_width=1,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )

        box = slide.shapes.add_textbox(
            Inches(CONTENT_LEFT + 0.2),
            Inches(top + 0.12),
            Inches(CONTENT_WIDTH - 0.4),
            Inches(height - 0.24),
        )
        tf = box.text_frame
        tf.word_wrap = True

        paragraphs: list[list[dict]] = [[]]
        for r in runs:
            parts = r["text"].split("\n")
            for part_idx, part in enumerate(parts):
                if part_idx > 0:
                    paragraphs.append([])
                if part:
                    paragraphs[-1].append({**r, "text": part})

        for p_idx, p_runs in enumerate(paragraphs):
            p = tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            if not p_runs:
                run = p.add_run()
                run.text = ""
                _set_font(run, name=CODE_FONT, size=10, color=CODE_DEFAULT)
            else:
                for r in p_runs:
                    run = p.add_run()
                    run.text = r["text"]
                    _set_font(
                        run,
                        name=CODE_FONT,
                        size=10,
                        bold=r.get("bold", False),
                        color=r.get("color", CODE_DEFAULT),
                    )

    # ── Tenets ───────────────────────────────────────────────────────────────

    def _add_tenets(self, slide, tenets: list[Tag], top: float):
        """Add tenet boxes to the slide."""
        num_tenets = len(tenets)
        if num_tenets >= 4:
            for i, tenet in enumerate(tenets):
                col = i % 2
                row = i // 2
                self._add_single_tenet(
                    slide, tenet, CONTENT_LEFT + col * 4.5, top + row * 1.0, width=4.2
                )
        else:
            for i, tenet in enumerate(tenets):
                self._add_single_tenet(
                    slide, tenet, CONTENT_LEFT, top + i * 1.0, width=CONTENT_WIDTH
                )

    def _add_single_tenet(
        self, slide, tenet: Tag, left: float, top: float, width: float
    ):
        """Add a single tenet box."""
        title_el = tenet.find(class_="tenet-title")
        text_el = tenet.find(class_="tenet-text")
        title = get_text(title_el) if title_el else ""
        text = get_text(text_el) if text_el else ""
        classes = tenet.get("class", [])
        accent_color = parse_color_from_class(classes) or MS_GREEN
        add_tenet(slide, title, text, left, top, width=width, accent_color=accent_color)

    # ── Versus comparison ────────────────────────────────────────────────────

    def _add_versus(self, slide, versus: Tag, top: float):
        """Add a versus comparison layout."""
        sides = versus.find_all(class_="versus-side")
        if len(sides) < 2:
            return

        for side_idx, side in enumerate(sides):
            base_left = CONTENT_LEFT if side_idx == 0 else 5.5

            title_el = side.find(class_="versus-title")
            title_h = 0.4
            if title_el:
                classes = title_el.get("class", [])
                color = parse_color_from_class(classes) or (
                    MS_ORANGE if side_idx == 0 else MS_GREEN
                )
                vs_title_text = get_text(title_el)
                title_h = _estimate_text_height(vs_title_text, 24, 4.0, bold=True)
                title_h = max(0.4, title_h)
                add_text_box(
                    slide,
                    vs_title_text,
                    left=base_left,
                    top=top,
                    width=4.0,
                    height=title_h,
                    font_size=24,
                    bold=True,
                    color=color,
                )

            items_list = side.find(class_="feature-list")
            if items_list:
                item_top = top + title_h + 0.1
                for item in items_list.find_all("li"):
                    text = get_text(item)
                    if "\u2713" in text:
                        color = MS_GREEN
                    elif "\u2717" in text:
                        color = MS_RED
                    else:
                        color = WHITE
                    item_h = _estimate_text_height(text, 14, 4.0)
                    item_h = max(0.30, item_h)
                    add_text_box(
                        slide,
                        text,
                        left=base_left,
                        top=item_top,
                        width=4.0,
                        height=item_h,
                        font_size=14,
                        color=color,
                    )
                    item_top += item_h

        # VS divider
        add_text_box(
            slide,
            "vs",
            left=4.5,
            top=top + 1.2,
            width=1.0,
            height=0.5,
            font_size=32,
            bold=True,
            color=GRAY_50,
            align=PP_ALIGN.CENTER,
        )

    # ── Tables ───────────────────────────────────────────────────────────────

    def _add_table(self, slide, table: Tag, top: float):
        """Add a data table to the slide."""
        rows = table.find_all("tr")
        if not rows:
            return

        # Pre-compute column widths
        sample_cells = rows[0].find_all(["th", "td"]) if rows else []
        num_cols = len(sample_cells)
        if num_cols == 0:
            return
        col_widths = [CONTENT_WIDTH / num_cols] * num_cols
        if num_cols == 3:
            col_widths = [2.5, 2.5, 3.4]
        elif num_cols == 2:
            col_widths = [4.2, 4.2]

        # Pre-compute per-row heights based on content
        row_heights: list[float] = []
        for row in rows:
            cells = row.find_all(["th", "td"])
            is_header = row.find("th") is not None
            fs = 12 if is_header else 11
            max_h = 0.32  # minimum row height
            for col_idx, cell in enumerate(cells):
                cw = col_widths[col_idx] if col_idx < len(col_widths) else 2.0
                cell_bold = is_header or col_idx == 0
                h = _estimate_text_height(get_text(cell), fs, cw, bold=cell_bold)
                max_h = max(max_h, h)
            row_heights.append(min(max_h, 0.8))  # cap individual rows

        current_row_top = top
        for row_idx, row in enumerate(rows):
            cells = row.find_all(["th", "td"])
            is_header = row.find("th") is not None
            r_num_cols = len(cells)
            if r_num_cols == 0:
                continue

            rh = row_heights[row_idx]
            left = CONTENT_LEFT
            for col_idx, cell in enumerate(cells):
                text = get_text(cell)
                width = col_widths[col_idx] if col_idx < len(col_widths) else 2.0

                if is_header:
                    color = self.accent_color
                    font_size = 12
                    bold = True
                else:
                    if col_idx == 0:
                        color = WHITE
                        bold = True
                    else:
                        color = GRAY_70
                        bold = False
                    font_size = 11

                    if "\u2713" in text:
                        color = MS_GREEN
                    elif "\u2717" in text:
                        color = MS_RED
                    elif "~" in text:
                        color = MS_ORANGE

                # Check for inline accent color (style attribute)
                style = cell.get("style", "")
                if "color:" in style and "var(--accent)" in style:
                    color = self.accent_color
                    bold = True

                add_text_box(
                    slide,
                    text,
                    left=left,
                    top=current_row_top,
                    width=width,
                    height=rh,
                    font_size=font_size,
                    bold=bold,
                    color=color,
                )
                left += width
            current_row_top += rh

    # ── Feature lists ────────────────────────────────────────────────────────

    def _add_feature_list(self, slide, feature_list: Tag, top: float):
        """Add a feature list to the slide."""
        items = feature_list.find_all("li")
        current_top = top
        for item in items:
            text = get_text(item)
            if "\u2713" in text:
                color = MS_GREEN
            elif "\u2717" in text:
                color = MS_RED
            else:
                color = WHITE
            item_h = max(0.4, _estimate_text_height(text, 16, CONTENT_WIDTH))
            add_text_box(
                slide,
                text,
                left=CONTENT_LEFT,
                top=current_top,
                width=CONTENT_WIDTH,
                height=item_h,
                font_size=16,
                color=color,
            )
            current_top += item_h

    # ── Stats grid/row ───────────────────────────────────────────────────────

    def _add_stats(self, slide, stat_container: Tag, top: float):
        """Add a stats grid or stat row to the slide.

        Handles both stat-grid/stat-row (stat, stat-number, stat-label)
        and velocity-grid (velocity-stat, velocity-number, velocity-label).
        """
        stats = stat_container.find_all(class_="stat") or stat_container.find_all(
            class_="velocity-stat"
        )
        num_stats = len(stats)
        if num_stats == 0:
            return

        width_per_stat = CONTENT_WIDTH / num_stats
        start_left = CONTENT_LEFT

        for i, stat in enumerate(stats):
            # Try all naming conventions (stat-* and velocity-*)
            number_el = (
                stat.find(class_="stat-number")
                or stat.find(class_="stat-value")
                or stat.find(class_="velocity-number")
            )
            label_el = stat.find(class_="stat-label") or stat.find(
                class_="velocity-label"
            )

            number = get_text(number_el) if number_el else ""
            label = get_text(label_el) if label_el else ""

            left = start_left + i * width_per_stat

            num_h = _estimate_text_height(number, 40, width_per_stat, bold=True)
            num_h = max(0.6, min(num_h, 1.2))
            add_text_box(
                slide,
                number,
                left=left,
                top=top,
                width=width_per_stat,
                height=num_h,
                font_size=40,
                bold=True,
                color=MS_CYAN,
                align=PP_ALIGN.CENTER,
            )
            label_h = _estimate_text_height(label, 12, width_per_stat)
            label_h = max(0.3, min(label_h, 0.6))
            add_text_box(
                slide,
                label,
                left=left,
                top=top + num_h,
                width=width_per_stat,
                height=label_h,
                font_size=12,
                color=GRAY_70,
                align=PP_ALIGN.CENTER,
            )

    # ── Quote ────────────────────────────────────────────────────────────────

    def _add_quote(self, slide, quote: Tag, top: float):
        """Add a quote to the slide."""
        text = get_text(quote)
        q_h = _estimate_text_height(f'"{text}"', 24, CONTENT_WIDTH)
        q_h = max(0.6, q_h)
        add_text_box(
            slide,
            f'"{text}"',
            left=CONTENT_LEFT,
            top=top,
            width=CONTENT_WIDTH,
            height=q_h,
            font_size=24,
            italic=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        attribution = quote.find_next_sibling(
            class_="quote-attribution"
        ) or quote.find_next_sibling(class_="quote-attr")
        if attribution:
            attr_text = get_text(attribution)
            attr_h = _estimate_text_height(attr_text, 14, CONTENT_WIDTH)
            attr_h = max(0.3, attr_h)
            add_text_box(
                slide,
                attr_text,
                left=CONTENT_LEFT,
                top=top + q_h,
                width=CONTENT_WIDTH,
                height=attr_h,
                font_size=14,
                color=GRAY_50,
                align=PP_ALIGN.CENTER,
            )

    # ── Convert & Save ───────────────────────────────────────────────────────

    def convert(self) -> Presentation:
        """Convert the HTML to a PowerPoint presentation."""
        slides = self.extract_slides()

        if not slides:
            warnings.warn(
                "No slides found in HTML. Check for <div class='slide'> or <section class='slide'> elements."
            )

        for i, slide_div in enumerate(slides):
            self.process_slide(slide_div, i + 1)

        if self.warnings:
            print(f"\nWarnings ({len(self.warnings)}):", file=sys.stderr)
            for w in self.warnings:
                print(f"  - {w}", file=sys.stderr)

        return self.prs

    def save(self, output_path: str):
        """Save the presentation to a file."""
        self.prs.save(output_path)


# ── CLI entry point ──────────────────────────────────────────────────────────


def main():
    parser = argparse.ArgumentParser(
        description="Convert Amplifier Stories HTML decks to PowerPoint presentations."
    )
    parser.add_argument("input", help="Input HTML file path")
    parser.add_argument(
        "output", nargs="?", help="Output PPTX file path (default: same name as input)"
    )

    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    if args.output:
        output_path = Path(args.output)
    else:
        output_path = input_path.with_suffix(".pptx")

    print(f"Converting: {input_path}")
    print(f"Output: {output_path}")

    html_content = input_path.read_text(encoding="utf-8")
    converter = HTMLToPPTXConverter(html_content)
    converter.convert()
    converter.save(str(output_path))

    num_slides = len(converter.prs.slides)
    print(f"Done! Created {output_path} ({num_slides} slides)")


if __name__ == "__main__":
    main()
