#!/usr/bin/env python3
"""
PPTX Text Overflow & Overlap Verifier (Layer 1 Detection)

Checks every text shape in a PPTX file to determine if the text
fits within the shape bounds. Uses per-character width tables for
accurate variable-width font estimation. Also detects overlapping shapes.

Usage:
    python tools/pptx_verify.py path/to/file.pptx
    python tools/pptx_verify.py path/to/directory/  # checks all .pptx files
"""

import math
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation

# ---------------------------------------------------------------------------
# Arial per-character width tables (fraction of em-size)
# ---------------------------------------------------------------------------
_ARIAL_REGULAR: dict[str, float] = {
    " ": 0.28,
    "!": 0.28,
    '"': 0.35,
    "#": 0.56,
    "$": 0.56,
    "%": 0.89,
    "&": 0.67,
    "'": 0.19,
    "(": 0.33,
    ")": 0.33,
    "*": 0.39,
    "+": 0.58,
    ",": 0.28,
    "-": 0.33,
    ".": 0.28,
    "/": 0.28,
    "0": 0.56,
    "1": 0.56,
    "2": 0.56,
    "3": 0.56,
    "4": 0.56,
    "5": 0.56,
    "6": 0.56,
    "7": 0.56,
    "8": 0.56,
    "9": 0.56,
    ":": 0.28,
    ";": 0.28,
    "<": 0.58,
    "=": 0.58,
    ">": 0.58,
    "?": 0.56,
    "@": 1.02,
    "A": 0.67,
    "B": 0.67,
    "C": 0.72,
    "D": 0.72,
    "E": 0.67,
    "F": 0.61,
    "G": 0.78,
    "H": 0.72,
    "I": 0.28,
    "J": 0.50,
    "K": 0.67,
    "L": 0.56,
    "M": 0.83,
    "N": 0.72,
    "O": 0.78,
    "P": 0.67,
    "Q": 0.78,
    "R": 0.72,
    "S": 0.67,
    "T": 0.61,
    "U": 0.72,
    "V": 0.67,
    "W": 0.94,
    "X": 0.67,
    "Y": 0.67,
    "Z": 0.61,
    "[": 0.28,
    "\\": 0.28,
    "]": 0.28,
    "^": 0.47,
    "_": 0.56,
    "`": 0.33,
    "a": 0.56,
    "b": 0.56,
    "c": 0.50,
    "d": 0.56,
    "e": 0.56,
    "f": 0.28,
    "g": 0.56,
    "h": 0.56,
    "i": 0.22,
    "j": 0.22,
    "k": 0.50,
    "l": 0.22,
    "m": 0.83,
    "n": 0.56,
    "o": 0.56,
    "p": 0.56,
    "q": 0.56,
    "r": 0.33,
    "s": 0.50,
    "t": 0.28,
    "u": 0.56,
    "v": 0.50,
    "w": 0.72,
    "x": 0.50,
    "y": 0.50,
    "z": 0.50,
    "{": 0.33,
    "|": 0.26,
    "}": 0.33,
    "~": 0.58,
}
_ARIAL_REGULAR_DEFAULT = 0.56
_ARIAL_BOLD_SCALE = 1.08

# Consolas is monospace -- single factor for all characters
_CONSOLAS_FACTOR = 0.60

# PowerPoint line height multiplier (single spacing = 1.0, but actual
# rendered height includes ascent + descent + internal leading ~ 1.2x)
LINE_HEIGHT_FACTOR = 1.2

# Default text frame margins in python-pptx (EMU)
DEFAULT_MARGIN_TOP = 45720  # 0.05"
DEFAULT_MARGIN_BOTTOM = 45720
DEFAULT_MARGIN_LEFT = 91440  # 0.1"
DEFAULT_MARGIN_RIGHT = 91440

EMU_PER_INCH = 914400
EMU_PER_PT = 12700

# Slide dimensions
SLIDE_WIDTH = 10.0  # inches
SLIDE_HEIGHT = 5.625  # inches (16:9 at 10" wide)


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------


@dataclass
class TextOverflow:
    """A detected text overflow issue."""

    slide_num: int
    shape_index: int
    text_preview: str
    font_size_pt: float
    font_name: str
    is_bold: bool
    shape_width: float  # inches
    shape_height: float  # inches
    available_height: float  # inches (after margins)
    needed_height: float  # inches
    overflow_inches: float  # how much it overflows by
    estimated_lines: int
    shape_top: float  # inches
    shape_bottom: float  # inches
    effective_bottom: float  # where text actually reaches


@dataclass
class ShapeOverlap:
    """A detected overlap between two shapes."""

    slide_num: int
    shape_a_index: int
    shape_b_index: int
    shape_a_text: str
    shape_b_text: str
    overlap_width: float  # inches
    overlap_height: float  # inches
    severity: str  # MINOR, MODERATE, SEVERE


@dataclass
class SlideReport:
    """Verification report for a single slide."""

    slide_num: int
    total_shapes: int
    text_shapes: int
    overflows: list[TextOverflow] = field(default_factory=list)
    overlaps: list[ShapeOverlap] = field(default_factory=list)

    @property
    def has_issues(self) -> bool:
        return len(self.overflows) > 0 or len(self.overlaps) > 0


@dataclass
class DeckReport:
    """Verification report for an entire deck."""

    path: str
    total_slides: int
    slides: list[SlideReport] = field(default_factory=list)

    @property
    def total_overflows(self) -> int:
        return sum(len(s.overflows) for s in self.slides)

    @property
    def total_overlaps(self) -> int:
        return sum(len(s.overlaps) for s in self.slides)

    @property
    def slides_with_issues(self) -> int:
        return sum(1 for s in self.slides if s.has_issues)


# ---------------------------------------------------------------------------
# Per-character width estimation
# ---------------------------------------------------------------------------


def _estimate_text_width_pt(
    text: str,
    font_size_pt: float,
    bold: bool = False,
    font_name: str = "Arial",
) -> float:
    """Estimate the rendered width of *text* in points.

    For Arial/Calibri uses per-character lookup from the width table.
    For Consolas (monospace) uses a single fixed factor.
    """
    if not text:
        return 0.0

    is_monospace = font_name.lower() in ("consolas", "courier new", "courier")

    if is_monospace:
        return len(text) * _CONSOLAS_FACTOR * font_size_pt

    # Variable-width (Arial / Calibri) -- per-character lookup
    total_em = 0.0
    for ch in text:
        total_em += _ARIAL_REGULAR.get(ch, _ARIAL_REGULAR_DEFAULT)

    width_pt = total_em * font_size_pt
    if bold:
        width_pt *= _ARIAL_BOLD_SCALE
    return width_pt


# ---------------------------------------------------------------------------
# Font info helpers
# ---------------------------------------------------------------------------


def _get_para_font_info(para) -> tuple[float, bool, str]:
    """Extract font size (pt), bold, and name from a paragraph."""
    font_size_pt = 14.0  # default
    is_bold = False
    font_name = "Arial"

    for run in para.runs:
        if run.font.size is not None:
            font_size_pt = run.font.size / EMU_PER_PT
        if run.font.bold:
            is_bold = True
        if run.font.name:
            font_name = run.font.name
        break  # use first run's properties

    return font_size_pt, is_bold, font_name


def _get_line_spacing(para) -> float:
    """Get effective line spacing multiplier from a paragraph."""
    if para.line_spacing is not None:
        if isinstance(para.line_spacing, float):
            return para.line_spacing
        # Could be in EMU (spacing points) -- convert
        if para.line_spacing > 10:
            return para.line_spacing / EMU_PER_PT / 12  # rough conversion
    return LINE_HEIGHT_FACTOR


# ---------------------------------------------------------------------------
# Height estimation (width-based wrapping)
# ---------------------------------------------------------------------------


def estimate_text_height(
    text: str,
    available_width_in: float,
    font_size_pt: float,
    font_name: str = "Arial",
    is_bold: bool = False,
    line_spacing: float = LINE_HEIGHT_FACTOR,
) -> tuple[float, int]:
    """Estimate height needed for text, and number of lines.

    Uses per-character width tables to compute rendered width, then
    determines how many wrap lines each paragraph needs.

    Returns:
        (height_in_inches, total_lines)
    """
    if not text.strip():
        return 0.0, 0

    available_width_pt = available_width_in * 72.0

    total_lines = 0.0
    paragraphs = text.split("\n")

    for para in paragraphs:
        stripped = para.strip()
        if not stripped:
            # Empty paragraph gets ~40% of a line height
            total_lines += 0.4
        else:
            rendered_width = _estimate_text_width_pt(
                stripped,
                font_size_pt,
                is_bold,
                font_name,
            )
            if rendered_width <= available_width_pt:
                total_lines += 1
            else:
                # Word-boundary wrapping is less efficient than char wrapping
                raw_lines = rendered_width / available_width_pt * 1.05
                total_lines += math.ceil(raw_lines)

    line_height_inches = font_size_pt / 72.0 * line_spacing

    return total_lines * line_height_inches, int(math.ceil(total_lines))


# ---------------------------------------------------------------------------
# Multi-paragraph height estimation
# ---------------------------------------------------------------------------


def _estimate_shape_text_height(
    tf,
    available_width_in: float,
) -> tuple[float, int, float, bool, str]:
    """Estimate total text height by iterating ALL paragraphs.

    Each paragraph may have different font properties. Sums up
    per-paragraph height estimates.

    Returns:
        (total_height_in, total_lines, dominant_font_size, dominant_bold,
         dominant_font_name)
    """
    total_height = 0.0
    total_lines = 0
    # Track the "dominant" font for reporting (largest font in shape)
    dominant_size = 0.0
    dominant_bold = False
    dominant_name = "Arial"

    for para in tf.paragraphs:
        font_size_pt, is_bold, font_name = _get_para_font_info(para)
        line_spacing = _get_line_spacing(para)

        # Track dominant (largest) font for the report
        if font_size_pt > dominant_size:
            dominant_size = font_size_pt
            dominant_bold = is_bold
            dominant_name = font_name

        para_text = para.text.strip()
        if not para_text:
            # Empty paragraph gets ~40% of a line height
            line_h = font_size_pt / 72.0 * line_spacing
            total_height += 0.4 * line_h
            continue

        # Compute rendered width for this paragraph
        available_width_pt = available_width_in * 72.0
        rendered_width = _estimate_text_width_pt(
            para_text,
            font_size_pt,
            is_bold,
            font_name,
        )

        if rendered_width <= available_width_pt:
            para_lines = 1
        else:
            raw_lines = rendered_width / available_width_pt * 1.05
            para_lines = math.ceil(raw_lines)

        line_h = font_size_pt / 72.0 * line_spacing
        total_height += para_lines * line_h
        total_lines += para_lines

    if dominant_size == 0.0:
        dominant_size = 14.0

    return total_height, total_lines, dominant_size, dominant_bold, dominant_name


# ---------------------------------------------------------------------------
# Overlap detection
# ---------------------------------------------------------------------------


def _shape_bbox(shape) -> tuple[float, float, float, float]:
    """Return bounding box (left, top, right, bottom) in inches."""
    left = shape.left / EMU_PER_INCH
    top = shape.top / EMU_PER_INCH
    right = left + shape.width / EMU_PER_INCH
    bottom = top + shape.height / EMU_PER_INCH
    return (left, top, right, bottom)


def _shape_preview(shape, max_len: int = 30) -> str:
    """Short text preview for a shape."""
    if shape.has_text_frame:
        txt = shape.text_frame.text.replace("\n", "|")[:max_len]
        return txt if txt.strip() else "(empty)"
    return "(no text)"


def detect_overlaps(slide, slide_num: int) -> list[ShapeOverlap]:
    """Detect overlapping shape bounding-boxes on a slide.

    Filters out trivially small overlaps (< 0.1" in both dimensions).
    Categorises severity by the larger overlap dimension:
        MINOR    < 0.2"
        MODERATE 0.2" -- 0.5"
        SEVERE   > 0.5"
    """
    shapes = list(slide.shapes)
    bboxes = [(i, _shape_bbox(s), s) for i, s in enumerate(shapes)]
    overlaps: list[ShapeOverlap] = []

    for ai in range(len(bboxes)):
        a_idx, (al, at, ar, ab), a_shape = bboxes[ai]
        a_text = (a_shape.text_frame.text.strip()
                  if a_shape.has_text_frame else "")
        for bi in range(ai + 1, len(bboxes)):
            b_idx, (bl, bt, br, bb), b_shape = bboxes[bi]
            b_text = (b_shape.text_frame.text.strip()
                      if b_shape.has_text_frame else "")

            # Skip intentional layering: background rects behind text.
            # Only report overlaps where BOTH shapes carry visible text.
            if not a_text or not b_text:
                continue

            # Intersection rectangle
            inter_left = max(al, bl)
            inter_top = max(at, bt)
            inter_right = min(ar, br)
            inter_bottom = min(ab, bb)

            overlap_w = inter_right - inter_left
            overlap_h = inter_bottom - inter_top

            if overlap_w <= 0 or overlap_h <= 0:
                continue  # no overlap

            # Filter trivially small overlaps.
            # Use the MIN dimension (the overlap *thickness*) to judge
            # significance.  A 4"×0.05" overlap between adjacent table
            # rows is thin/benign, even though the width is large.
            min_dim = min(overlap_w, overlap_h)
            if min_dim < 0.10:
                continue  # barely touching — skip

            if min_dim > 0.50:
                severity = "SEVERE"
            elif min_dim > 0.20:
                severity = "MODERATE"
            else:
                severity = "MINOR"

            overlaps.append(
                ShapeOverlap(
                    slide_num=slide_num,
                    shape_a_index=a_idx,
                    shape_b_index=b_idx,
                    shape_a_text=_shape_preview(a_shape),
                    shape_b_text=_shape_preview(b_shape),
                    overlap_width=overlap_w,
                    overlap_height=overlap_h,
                    severity=severity,
                )
            )

    return overlaps


# ---------------------------------------------------------------------------
# Shape / slide / deck verification
# ---------------------------------------------------------------------------


def verify_shape(shape, slide_num: int, shape_idx: int) -> TextOverflow | None:
    """Check if text fits within a shape. Returns overflow info or None."""
    if not shape.has_text_frame:
        return None

    text = shape.text_frame.text
    if not text.strip():
        return None

    tf = shape.text_frame

    # Shape dimensions
    shape_w = shape.width / EMU_PER_INCH
    shape_h = shape.height / EMU_PER_INCH
    shape_top = shape.top / EMU_PER_INCH

    # Text frame margins
    m_top = (
        tf.margin_top if tf.margin_top is not None else DEFAULT_MARGIN_TOP
    ) / EMU_PER_INCH
    m_bottom = (
        tf.margin_bottom if tf.margin_bottom is not None else DEFAULT_MARGIN_BOTTOM
    ) / EMU_PER_INCH
    m_left = (
        tf.margin_left if tf.margin_left is not None else DEFAULT_MARGIN_LEFT
    ) / EMU_PER_INCH
    m_right = (
        tf.margin_right if tf.margin_right is not None else DEFAULT_MARGIN_RIGHT
    ) / EMU_PER_INCH

    avail_w = shape_w - m_left - m_right
    avail_h = shape_h - m_top - m_bottom

    if avail_w <= 0 or avail_h <= 0:
        return None  # degenerate shape, skip

    # Check if word wrap is enabled (affects whether text wraps or extends)
    word_wrap = tf.word_wrap
    if word_wrap is False:
        # No wrapping -- text extends horizontally, only newlines create new lines
        # Still use multi-paragraph estimation for mixed fonts
        needed_h, num_lines, font_size_pt, is_bold, font_name = (
            _estimate_shape_text_height(tf, avail_w)
        )
        # Override: for no-wrap, count only explicit newlines
        num_lines = max(1, len(text.split("\n")))
        line_height = font_size_pt / 72.0 * LINE_HEIGHT_FACTOR
        needed_h = num_lines * line_height
    else:
        # Word wrap enabled (or default) -- multi-paragraph estimation
        needed_h, num_lines, font_size_pt, is_bold, font_name = (
            _estimate_shape_text_height(tf, avail_w)
        )

    overflow = needed_h - avail_h

    # Threshold: flag if text needs more than 15% extra height
    # (small overflows may be absorbed by PowerPoint's text compression)
    if overflow > avail_h * 0.15 and overflow > 0.05:
        return TextOverflow(
            slide_num=slide_num,
            shape_index=shape_idx,
            text_preview=text.replace("\n", "|")[:60],
            font_size_pt=font_size_pt,
            font_name=font_name,
            is_bold=is_bold,
            shape_width=shape_w,
            shape_height=shape_h,
            available_height=avail_h,
            needed_height=needed_h,
            overflow_inches=overflow,
            estimated_lines=num_lines,
            shape_top=shape_top,
            shape_bottom=shape_top + shape_h,
            effective_bottom=shape_top + m_top + needed_h + m_bottom,
        )

    return None


def verify_slide(slide, slide_num: int) -> SlideReport:
    """Verify all text shapes in a slide and detect overlaps."""
    shapes = list(slide.shapes)
    text_shapes = [s for s in shapes if s.has_text_frame and s.text_frame.text.strip()]

    report = SlideReport(
        slide_num=slide_num,
        total_shapes=len(shapes),
        text_shapes=len(text_shapes),
    )

    for idx, shape in enumerate(shapes):
        overflow = verify_shape(shape, slide_num, idx)
        if overflow:
            report.overflows.append(overflow)

    # Overlap detection
    report.overlaps = detect_overlaps(slide, slide_num)

    return report


def verify_deck(path: str) -> DeckReport:
    """Verify all slides in a PPTX deck."""
    prs = Presentation(path)
    slides = list(prs.slides)

    report = DeckReport(path=path, total_slides=len(slides))

    for i, slide in enumerate(slides):
        slide_report = verify_slide(slide, i + 1)
        report.slides.append(slide_report)

    return report


# ---------------------------------------------------------------------------
# Report formatting
# ---------------------------------------------------------------------------


def format_report(report: DeckReport, verbose: bool = False) -> str:
    """Format a deck report as human-readable text."""
    lines = []
    name = Path(report.path).stem
    lines.append(f"\n{'=' * 70}")
    lines.append(f"{name}")
    lines.append(f"{'=' * 70}")

    if report.total_overflows == 0 and report.total_overlaps == 0:
        lines.append("  ALL CLEAN - no text overflow or overlap detected")
        return "\n".join(lines)

    parts = []
    if report.total_overflows:
        parts.append(f"{report.total_overflows} overflows")
    if report.total_overlaps:
        parts.append(f"{report.total_overlaps} overlaps")
    lines.append(
        f"  {', '.join(parts)} across "
        f"{report.slides_with_issues}/{report.total_slides} slides"
    )

    for sr in report.slides:
        if not sr.has_issues:
            if verbose:
                lines.append(f"\n  Slide {sr.slide_num}: CLEAN")
            continue

        issue_parts = []
        if sr.overflows:
            issue_parts.append(f"{len(sr.overflows)} overflow(s)")
        if sr.overlaps:
            issue_parts.append(f"{len(sr.overlaps)} overlap(s)")
        lines.append(f"\n  Slide {sr.slide_num}: {', '.join(issue_parts)}")

        for ov in sr.overflows:
            severity = (
                "SEVERE"
                if ov.overflow_inches > 0.5
                else "MODERATE"
                if ov.overflow_inches > 0.2
                else "MINOR"
            )
            lines.append(
                f"    [{severity}] {ov.font_size_pt:.0f}pt {ov.font_name}"
                f"{'(B)' if ov.is_bold else ''}"
                f' in {ov.shape_width:.1f}"x{ov.shape_height:.2f}"'
                f' needs {ov.needed_height:.2f}"'
                f' (overflow: {ov.overflow_inches:.2f}")'
            )
            lines.append(f'      "{ov.text_preview}"')
            if ov.effective_bottom > SLIDE_HEIGHT:
                lines.append(
                    f'      ^^ RUNS OFF SLIDE (bottom={ov.effective_bottom:.2f}")'
                )

        for ol in sr.overlaps:
            lines.append(
                f"    [OVERLAP-{ol.severity}] shapes {ol.shape_a_index} & "
                f"{ol.shape_b_index}: "
                f'{ol.overlap_width:.2f}"x{ol.overlap_height:.2f}" overlap'
            )
            lines.append(f'      A: "{ol.shape_a_text}"  B: "{ol.shape_b_text}"')

    return "\n".join(lines)


def main():
    """CLI entry point."""
    if len(sys.argv) < 2:
        print("Usage: python tools/pptx_verify.py <file_or_dir> [--verbose]")
        sys.exit(1)

    path = Path(sys.argv[1])
    verbose = "--verbose" in sys.argv

    if path.is_dir():
        pptx_files = sorted(path.glob("*.pptx"))
    else:
        pptx_files = [path]

    if not pptx_files:
        print(f"No .pptx files found at {path}")
        sys.exit(1)

    total_issues = 0
    total_overlaps = 0
    total_slides = 0
    total_clean = 0

    for pptx_file in pptx_files:
        report = verify_deck(str(pptx_file))
        print(format_report(report, verbose))
        total_issues += report.total_overflows
        total_overlaps += report.total_overlaps
        total_slides += report.total_slides
        total_clean += report.total_slides - report.slides_with_issues

    print(f"\n{'=' * 70}")
    print(
        f"SUMMARY: {total_issues} overflow issues, {total_overlaps} overlaps across {total_slides} slides"
    )
    print(f"  Clean slides: {total_clean}/{total_slides}")
    print(f"{'=' * 70}")


if __name__ == "__main__":
    main()
